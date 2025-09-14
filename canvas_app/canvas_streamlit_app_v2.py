#!/usr/bin/env python3
"""
Canvas LMS Streamlit App V2

A web interface for Canvas LMS that focuses on accessible data:
- Assignments and due dates
- Course modules and content
- Course information
- User data
- Assignment submissions (if available)
"""

import streamlit as st
import os
import json
import requests
import time
from pathlib import Path
from datetime import datetime
import zipfile
import io
from typing import List, Dict, Any, Optional, Tuple
import logging
import pandas as pd
import re
from urllib.parse import urlparse, unquote, parse_qs
import glob
import tempfile
import sys
import numpy as np
import hashlib
import html
import urllib.parse
import subprocess

from canvas_client import CanvasClient
from canvas_config import CanvasConfig
# RAG Demo imports
try:
    # Add send_to_friend directory to path for imports
    send_to_friend_path = os.path.join(os.path.dirname(__file__), 'send_to_friend')
    if send_to_friend_path not in sys.path:
        sys.path.insert(0, send_to_friend_path)
    
    # Import from correct paths based on directory structure
    # TextExtractionPipeline removed - using built-in method instead
    from src.utils.ingest import ingest_documents, semantic_search
    from src.utils.retrieval import mmr_retrieve
    from src.utils.rag_llm import answer_with_context
    from src.utils.memory import ConversationMemory
    from src.utils.learning_mode import extract_topics_from_notes, recommend_youtube, build_topic_context
    from src.utils.web_search import recommend_articles_ddg, recommend_youtube_ddg
    from src.utils.notes import generate_notes_from_text, iter_generate_notes_from_texts
    from src.utils.notes_ingest import ingest_notes_sections
except ImportError as e:
    st.error(f"RAG Demo dependencies not found: {e}")
    st.error("Please ensure the send_to_friend directory and its dependencies are available.")
    # Continue without RAG features
    RAG_AVAILABLE = False
else:
    RAG_AVAILABLE = True

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# RAG Demo Configuration
if RAG_AVAILABLE:
    TOP_K = 5
    EMBED_MODEL = "all-MiniLM-L6-v2"
    
    def _store_dir():
        return os.path.join("data", "vector_store")
    
    def _meta_path():
        return os.path.join("data", "notes_index", "metadata.json")
    
    # Initialize memory
    mem = ConversationMemory()

# Check for Navdeep components availability
try:
    navdeep_path = "/Users/praveenbhandari/sf-vibe copy/navdeep/src"
    if navdeep_path not in sys.path:
        sys.path.append(navdeep_path)
    # TextExtractionPipeline removed - using built-in method instead
    NAVDEEP_AVAILABLE = True
except ImportError as e:
    NAVDEEP_AVAILABLE = False
except Exception as e:
    NAVDEEP_AVAILABLE = False

# Page configuration
st.set_page_config(
    page_title="Canvas LMS Explorer + RAG Demo",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS
st.markdown("""
<style>
    .main-header {
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        padding: 1rem;
        border-radius: 10px;
        color: white;
        text-align: center;
        margin-bottom: 2rem;
    }
    .metric-card {
        background: white;
        padding: 1rem;
        border-radius: 8px;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        border-left: 4px solid #667eea;
    }
    .status-badge {
        padding: 0.25rem 0.5rem;
        border-radius: 12px;
        font-size: 0.8rem;
        font-weight: bold;
    }
    .status-completed { background-color: #d4edda; color: #155724; }
    .status-pending { background-color: #fff3cd; color: #856404; }
    .status-in-progress { background-color: #cce5ff; color: #004085; }
    .workflow-step {
        padding: 1rem;
        margin: 0.5rem 0;
        border-radius: 8px;
        border-left: 4px solid #28a745;
    }
    .sidebar .sidebar-content {
        background-color: #f8f9fa;
    }
    .chat-message {
        padding: 1rem;
        margin: 0.5rem 0;
        border-radius: 8px;
        background-color: #f8f9fa;
    }
    .user-message {
        background-color: #e3f2fd;
        border-left: 4px solid #2196f3;
    }
    .assistant-message {
        background-color: #f3e5f5;
        border-left: 4px solid #9c27b0;
    }
    .course-card {
        border: 1px solid #ddd;
        border-radius: 8px;
        padding: 1rem;
        margin: 0.5rem 0;
        background-color: #f9f9f9;
    }
    .assignment-card {
        border-left: 4px solid #4CAF50;
        padding: 1rem;
        margin: 0.5rem 0;
        background-color: #f8f9fa;
    }
    .module-card {
        border-left: 4px solid #2196F3;
        padding: 1rem;
        margin: 0.5rem 0;
        background-color: #f8f9fa;
    }
    .status-success {
        color: #28a745;
        font-weight: bold;
    }
    .status-error {
        color: #dc3545;
        font-weight: bold;
    }
    .due-date {
        color: #ff6b35;
        font-weight: bold;
    }
</style>
""", unsafe_allow_html=True)

class CanvasFileDownloader:
    """Download files from Canvas course JSON exports"""
    
    def __init__(self, api_token: str = None, base_url: str = None, output_dir: str = "downloads", dry_run: bool = False):
        """
        Initialize the downloader
        
        Args:
            api_token: Canvas API token (optional, for authenticated requests)
            base_url: Canvas base URL (optional)
            output_dir: Directory to save downloaded files
            dry_run: If True, only show what would be downloaded without downloading
        """
        self.api_token = api_token
        self.base_url = base_url
        self.output_dir = Path(output_dir)
        self.dry_run = dry_run
        self.session = requests.Session()
        
        # Set up session headers if token provided
        if self.api_token:
            self.session.headers.update({
                'Authorization': f'Bearer {self.api_token}',
                'User-Agent': 'Canvas-File-Downloader/1.0'
            })
        
        # Create output directory
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Statistics
        self.stats = {
            'total_files': 0,
            'downloaded': 0,
            'skipped': 0,
            'errors': 0,
            'external_urls': 0
        }
    
    def sanitize_filename(self, filename: str) -> str:
        """Sanitize filename for safe file system usage"""
        # Remove or replace invalid characters
        filename = re.sub(r'[<>:"/\\|?*]', '_', filename)
        # Remove leading/trailing spaces and dots
        filename = filename.strip(' .')
        # Limit length
        if len(filename) > 200:
            name, ext = os.path.splitext(filename)
            filename = name[:200-len(ext)] + ext
        return filename or 'unnamed_file'
    
    def get_filename_from_url(self, url: str, default_name: str = None) -> str:
        """Extract filename from URL"""
        try:
            parsed = urlparse(url)
            filename = os.path.basename(unquote(parsed.path))
            if filename and '.' in filename:
                return self.sanitize_filename(filename)
        except Exception:
            pass
        
        return self.sanitize_filename(default_name or 'downloaded_file')
    
    def download_file(self, url: str, filename: str, course_dir: Path, progress_callback=None) -> bool:
        """Download a single file with improved error handling
        
        Args:
            url: URL to download from
            filename: Suggested filename
            course_dir: Directory to save the file
            progress_callback: Optional callback for progress updates
            
        Returns:
            True if successful, False otherwise
        """
        try:
            # Create course directory if it doesn't exist
            course_dir.mkdir(parents=True, exist_ok=True)
            
            # Get filename from URL if not provided
            if not filename or filename == 'unknown_file':
                filename = self.get_filename_from_url(url)
            
            # Sanitize filename
            safe_filename = self.sanitize_filename(filename)
            filepath = course_dir / safe_filename
            
            # Skip if file already exists
            if filepath.exists():
                if progress_callback:
                    progress_callback(f"File already exists, skipping: {safe_filename}")
                self.stats['skipped'] += 1
                return True
            
            if progress_callback:
                progress_callback(f"{'[DRY RUN] Would download' if self.dry_run else 'Downloading'}: {safe_filename}")
            
            # In dry-run mode, just simulate the download
            if self.dry_run:
                self.stats['downloaded'] += 1
                return True
            
            # Download with timeout and streaming
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
            response = self.session.get(url, headers=headers, timeout=60, stream=True)
            response.raise_for_status()
            
            # Check content type for potential issues
            content_type = response.headers.get('content-type', '').lower()
            if 'text/html' in content_type and not filename.endswith('.html'):
                if progress_callback:
                    progress_callback(f"Warning: Received HTML content for {filename}, might be an error page")
            
            # Write file in chunks with verification
            temp_filepath = filepath.with_suffix(filepath.suffix + '.tmp')
            total_size = 0
            
            with open(temp_filepath, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    if chunk:
                        f.write(chunk)
                        total_size += len(chunk)
            
            # Verify download completed successfully
            if total_size == 0:
                if progress_callback:
                    progress_callback(f"Error: Downloaded file is empty: {safe_filename}")
                temp_filepath.unlink()
                self.stats['errors'] += 1
                return False
            
            # Move temp file to final location
            temp_filepath.rename(filepath)
            
            if progress_callback:
                progress_callback(f"Downloaded: {safe_filename} ({total_size} bytes)")
            self.stats['downloaded'] += 1
            return True
            
        except requests.exceptions.RequestException as e:
            if progress_callback:
                progress_callback(f"Network error downloading {filename}: {e}")
            self.stats['errors'] += 1
            return False
        except Exception as e:
            if progress_callback:
                progress_callback(f"Error downloading {filename}: {e}")
            self.stats['errors'] += 1
            return False
    
    def process_module_items(self, items: List[Dict], course_dir: Path, progress_callback=None) -> None:
        """Process module items and download content"""
        for item in items:
            self.stats['total_files'] += 1
            
            title = item.get('title', 'Untitled')
            item_type = item.get('type', 'Unknown')
            
            # Process different types of URLs
            if 'url' in item and item['url']:
                url = item['url']
                
                if item_type == 'File' or 'files' in url:
                    # Canvas file
                    filename = item.get('display_name', title)
                    self.download_file(url, filename, course_dir, progress_callback)
            
            elif 'external_url' in item and item['external_url']:
                url = item['external_url']
                
                if 'drive.google.com' in url:
                    # Google Drive file
                    self.process_google_drive_url(url, title, course_dir, progress_callback)
    
    def process_google_drive_url(self, url: str, title: str, course_dir: Path, progress_callback=None) -> bool:
        """Process Google Drive URL and attempt direct download"""
        try:
            # Convert Google Drive view URL to download URL
            if '/file/d/' in url and '/view' in url:
                file_id = url.split('/file/d/')[1].split('/')[0]
                download_url = f"https://drive.google.com/uc?export=download&id={file_id}"
                
                if progress_callback:
                    progress_callback(f"Converting Google Drive URL for: {title}")
                return self.download_file(download_url, f"{title}.pdf", course_dir, progress_callback)
            else:
                # Save as link if can't convert
                return self.process_external_url(url, title, course_dir, progress_callback)
                
        except Exception as e:
            if progress_callback:
                progress_callback(f"Error processing Google Drive URL {title}: {e}")
            return self.process_external_url(url, title, course_dir, progress_callback)
    
    def process_external_url(self, url: str, title: str, course_dir: Path, progress_callback=None) -> bool:
        """Process external URL (save link info)"""
        try:
            # For external URLs, save a link file
            safe_filename = self.sanitize_filename(f"{title}_link.txt")
            filepath = course_dir / safe_filename
            
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write(f"Title: {title}\n")
                f.write(f"URL: {url}\n")
                f.write(f"Saved: {datetime.now().isoformat()}\n")
            
            if progress_callback:
                progress_callback(f"Saved external link: {safe_filename}")
            self.stats['external_urls'] += 1
            return True
            
        except Exception as e:
            if progress_callback:
                progress_callback(f"Error saving external URL {title}: {e}")
            self.stats['errors'] += 1
            return False
    
    def process_course_files(self, files: List[Dict], course_dir: Path, progress_callback=None) -> None:
        """Process course files list"""
        files_dir = course_dir / "course_files"
        
        for file_info in files:
            self.stats['total_files'] += 1
            
            filename = file_info.get('display_name', file_info.get('filename', 'unknown_file'))
            file_url = file_info.get('url')
            
            if file_url:
                self.download_file(file_url, filename, files_dir, progress_callback)
    
    def process_json_export(self, json_file: str, progress_callback=None) -> None:
        """Process a single JSON export file"""
        try:
            if progress_callback:
                progress_callback(f"Processing JSON export: {json_file}")
            
            with open(json_file, 'r', encoding='utf-8') as f:
                course_data = json.load(f)
            
            # Handle both old and new export formats
            if 'course_data' in course_data:
                # New format with metadata
                actual_data = course_data['course_data']
                course_info = actual_data.get('course_info', {})
            else:
                # Old format
                actual_data = course_data
                course_info = actual_data.get('course_info', {})
            
            course_name = course_info.get('name', f"Course_{actual_data.get('course_id', 'Unknown')}")
            course_id = actual_data.get('course_id', 'unknown')
            
            # Create course directory
            safe_course_name = self.sanitize_filename(course_name)
            course_dir = self.output_dir / f"course_{course_id}_{safe_course_name}"
            
            if progress_callback:
                progress_callback(f"Processing course: {course_name} (ID: {course_id})")
            
            # Process modules and their items
            modules = actual_data.get('modules', [])
            for module in modules:
                module_name = module.get('name', 'Unnamed Module')
                if progress_callback:
                    progress_callback(f"Processing module: {module_name}")
                
                items = module.get('items', [])
                if items:
                    module_dir = course_dir / self.sanitize_filename(module_name)
                    self.process_module_items(items, module_dir, progress_callback)
            
            # Process course files
            files = actual_data.get('files', [])
            if files:
                if progress_callback:
                    progress_callback(f"Processing {len(files)} course files")
                self.process_course_files(files, course_dir, progress_callback)
            
            if progress_callback:
                progress_callback(f"Completed processing: {course_name}")
            
        except Exception as e:
            if progress_callback:
                progress_callback(f"Error processing JSON file {json_file}: {e}")
            self.stats['errors'] += 1


class CanvasCourseExplorer:
    """Canvas LMS Course Explorer with focus on accessible data"""
    
    def __init__(self):
        self.client = None
        self.config_manager = CanvasConfig()
        self.export_dir = Path("exports")
        self.export_dir.mkdir(exist_ok=True)
        self.download_dir = Path("downloads")
        self.download_dir.mkdir(exist_ok=True)
    
    def initialize_session_state(self):
        """Initialize session state variables"""
        # Canvas LMS session state
        if 'authenticated' not in st.session_state:
            st.session_state.authenticated = False
        if 'user_info' not in st.session_state:
            st.session_state.user_info = None
        if 'courses' not in st.session_state:
            st.session_state.courses = []
        if 'selected_course' not in st.session_state:
            st.session_state.selected_course = None
        if 'course_data' not in st.session_state:
            st.session_state.course_data = {}
        if 'client' not in st.session_state:
            st.session_state.client = None
        # For backward compatibility, also set canvas_client
        if 'canvas_client' not in st.session_state:
            st.session_state.canvas_client = None
        # Sync canvas_client with client when client is available
        if st.session_state.client and not st.session_state.canvas_client:
            st.session_state.canvas_client = st.session_state.client
            
        # RAG Demo session state
        if RAG_AVAILABLE:
            if 'uploaded_files' not in st.session_state:
                st.session_state.uploaded_files = []
            if 'notes_generated' not in st.session_state:
                st.session_state.notes_generated = False
            if 'chat_history' not in st.session_state:
                st.session_state.chat_history = []
            if 'current_notes' not in st.session_state:
                st.session_state.current_notes = ""
            if 'extracted_topics' not in st.session_state:
                st.session_state.extracted_topics = []
            if 'rag_selected_files' not in st.session_state:
                st.session_state.rag_selected_files = []
    
    def render_header(self):
        """Render the main header"""
        st.markdown("""
        <div class="main-header">
            <h1 style="margin: 0; font-size: 2.5rem;">🎓 Canvas LMS Explorer + RAG Demo</h1>
            <p style="margin: 10px 0 0 0; opacity: 0.9; font-size: 1.1rem;">Explore your courses and enhance learning with AI</p>
        </div>
        """, unsafe_allow_html=True)
    
    def render_login_section(self):
        """Render the login section"""
        st.sidebar.header("🔐 Authentication")
        
        # Check if already authenticated
        if st.session_state.authenticated:
            st.sidebar.success("✅ Authenticated")
            if st.session_state.user_info:
                st.sidebar.write(f"**User:** {st.session_state.user_info.get('name', 'Unknown')}")
                st.sidebar.write(f"**ID:** {st.session_state.user_info.get('id', 'Unknown')}")
            
            # RAG Demo status indicator
            st.sidebar.markdown("---")
            if RAG_AVAILABLE:
                st.sidebar.success("🤖 RAG Demo: Available")
            else:
                st.sidebar.error("🤖 RAG Demo: Not Available")
            
            if st.sidebar.button("🚪 Logout", type="secondary"):
                self.logout()
            return True
        
        # Login form
        with st.sidebar.form("login_form"):
            st.write("Enter your Canvas credentials:")
            
            canvas_url = st.text_input(
                "Canvas URL",
                value=self.config_manager.get('base_url', ''),
                placeholder="https://your-school.instructure.com",
                help="Your Canvas instance URL"
            )
            
            api_token = st.text_input(
                "API Token",
                type="password",
                placeholder="Enter your Canvas API token",
                help="Get this from Account → Settings → Approved Integrations"
            )
            
            submitted = st.form_submit_button("🔑 Login", type="primary")
            
            if submitted:
                if canvas_url and api_token:
                    self.login(canvas_url, api_token)
                else:
                    st.error("Please enter both Canvas URL and API token")
        
        # Instructions
        with st.sidebar.expander("ℹ️ How to get API token"):
            st.markdown("""
            1. Log into your Canvas account
            2. Go to **Account** → **Settings**
            3. Scroll down to **Approved Integrations**
            4. Click **+ New Access Token**
            5. Give it a name and generate
            6. Copy the token and paste it here
            """)
        
        return False
    
    def login(self, canvas_url: str, api_token: str):
        """Handle login process"""
        try:
            with st.spinner("Authenticating..."):
                # Create client with provided credentials
                self.client = CanvasClient(base_url=canvas_url, api_token=api_token)
                
                # Validate credentials
                user_info = self.client.validate_credentials()
                
                # Save configuration
                self.config_manager.save_config(
                    base_url=canvas_url,
                    api_token=api_token
                )
                
                # Update session state
                st.session_state.authenticated = True
                st.session_state.user_info = user_info
                st.session_state.client = self.client
                st.session_state.canvas_client = self.client  # For backward compatibility
                
                st.success(f"✅ Successfully authenticated as {user_info.get('name')}")
                st.rerun()
                
        except Exception as e:
            st.error(f"❌ Authentication failed: {str(e)}")
            logger.error(f"Login error: {e}")
    
    def logout(self):
        """Handle logout process"""
        st.session_state.authenticated = False
        st.session_state.user_info = None
        st.session_state.courses = []
        st.session_state.selected_course = None
        st.session_state.course_data = {}
        st.session_state.client = None
        st.session_state.canvas_client = None  # For backward compatibility
        st.rerun()
    
    def load_courses(self):
        """Load courses from Canvas"""
        if not st.session_state.authenticated or not st.session_state.client:
            return
        
        try:
            with st.spinner("Loading courses..."):
                courses = st.session_state.client.get_courses()
                st.session_state.courses = courses
                st.success(f"✅ Loaded {len(courses)} courses")
        except Exception as e:
            st.error(f"❌ Failed to load courses: {str(e)}")
            logger.error(f"Load courses error: {e}")
    
    def _load_local_course_files(self, course_id: int):
        """Load course files from local downloads directory"""
        try:
            downloads_dir = Path("downloads")
            if not downloads_dir.exists():
                return []
            
            # Look for course-specific folders
            course_patterns = [
                f"course_{course_id}_*",
                f"*{course_id}*",
                "*"  # Fallback to all folders
            ]
            
            local_files = []
            supported_extensions = ['.pdf', '.docx', '.doc', '.txt', '.pptx', '.ppsx', '.mp4', '.mp3']
            
            for pattern in course_patterns:
                course_dirs = list(downloads_dir.glob(pattern))
                if course_dirs:
                    for course_dir in course_dirs:
                        if course_dir.is_dir():
                            for ext in supported_extensions:
                                files = list(course_dir.rglob(f'*{ext}'))
                                for file_path in files:
                                    # Create file info similar to Canvas API format
                                    file_info = {
                                        'id': hash(str(file_path)),  # Generate unique ID
                                        'display_name': file_path.name,
                                        'filename': file_path.name,
                                        'size': file_path.stat().st_size,
                                        'content-type': self._get_content_type(file_path),
                                        'url': str(file_path),
                                        'created_at': datetime.fromtimestamp(file_path.stat().st_ctime).isoformat(),
                                        'updated_at': datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
                                        'local_path': str(file_path),
                                        'is_local': True
                                    }
                                    local_files.append(file_info)
                    break  # Found course files, no need to check other patterns
            
            logger.info(f"Found {len(local_files)} local files for course {course_id}")
            return local_files
            
        except Exception as e:
            logger.warning(f"Could not load local files: {e}")
            return []
    
    def _get_content_type(self, file_path: Path):
        """Get content type based on file extension"""
        ext = file_path.suffix.lower()
        content_types = {
            '.pdf': 'application/pdf',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.doc': 'application/msword',
            '.txt': 'text/plain',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.ppsx': 'application/vnd.ms-powerpoint',
            '.mp4': 'video/mp4',
            '.mp3': 'audio/mpeg'
        }
        return content_types.get(ext, 'application/octet-stream')
    
    # RAG and Memory Management Classes
    class MemoryStore:
        """Simple JSON-backed memory store for short-term and long-term memory."""
        
        def __init__(self, root_dir: str = "data/memory") -> None:
            self.root_dir = root_dir
            os.makedirs(self.root_dir, exist_ok=True)

        def _lt_path(self, profile_id: str) -> str:
            fname = f"long_term_{profile_id}.json"
            return os.path.join(self.root_dir, fname)

        def load_long_term(self, profile_id: str = "default") -> Dict[str, Any]:
            path = self._lt_path(profile_id)
            if not os.path.exists(path):
                return {"facts": [], "topic_progress": {}}
            try:
                with open(path, "r", encoding="utf-8") as f:
                    return json.load(f)
            except Exception:
                return {"facts": [], "topic_progress": {}}

        def save_long_term(self, data: Dict[str, Any], profile_id: str = "default") -> None:
            path = self._lt_path(profile_id)
            with open(path, "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=2)

        @staticmethod
        def format_recent(messages: List[Dict[str, str]], window: int = 6) -> str:
            recent = messages[-window:]
            lines: List[str] = []
            for m in recent:
                role = m.get("role", "user")
                content = m.get("content", "")
                lines.append(f"{role}: {content}")
            return "\n".join(lines)

        def summarize_and_store_long_term(
            self,
            messages: List[Dict[str, str]],
            profile_id: str = "default",
            min_turns: int = 8,
            max_turns: int = 16,
        ) -> Optional[List[str]]:
            """Summarize recent conversation into persistent facts and merge them into long-term memory."""
            if len(messages) < min_turns:
                return None
            context_text = self.format_recent(messages, window=max_turns)
            contexts = [{
                "source": "conversation",
                "chunk_index": 0,
                "text": context_text,
            }]
            prompt = (
                "From the conversation, extract persistent student facts/preferences/goals and any difficulties. "
                "Output 3-6 concise bullet points. Avoid ephemeral content."
            )
            try:
                resp = self._answer_with_context(prompt, contexts)
            except Exception:
                return None
            # Parse bullets
            facts = []
            for line in resp.splitlines():
                l = line.strip().lstrip("- ")
                if l:
                    facts.append(l)
            if not facts:
                return None
            lt = self.load_long_term(profile_id)
            existing = set(lt.get("facts", []))
            changed = False
            for f in facts:
                if f not in existing:
                    existing.add(f)
                    changed = True
            if changed:
                lt["facts"] = list(existing)
                self.save_long_term(lt, profile_id)
                return facts
            return None

        def memory_contexts(
            self,
            messages: List[Dict[str, str]],
            profile_id: str = "default",
            short_window: int = 6,
        ) -> List[Dict[str, Any]]:
            """Build contexts representing short-term and long-term memory."""
            contexts: List[Dict[str, Any]] = []
            st_text = self.format_recent(messages, window=short_window)
            if st_text:
                contexts.append({"source": "memory:short_term", "chunk_index": 0, "text": st_text})
            lt = self.load_long_term(profile_id)
            facts = lt.get("facts", [])
            if facts:
                contexts.append({"source": "memory:long_term", "chunk_index": 0, "text": "\n".join(f"- {f}" for f in facts)})
            return contexts

        def _answer_with_context(self, query: str, contexts: List[Dict[str, Any]]) -> str:
            """Simple context-based answer generation"""
            context_text = "\n\n".join([f"[source: {c.get('source')}]\n{c.get('text', '')}" for c in contexts])
            return f"Based on the context:\n\n{context_text}\n\nQuery: {query}\n\nAnswer: [This is a simplified response. For full RAG functionality, configure external LLM services.]"

    # Web Search and Learning Mode Utilities
    def recommend_articles_ddg(self, topic: str, limit: int = 3) -> List[Dict[str, str]]:
        """Fetch simple article links from DuckDuckGo HTML endpoint."""
        q = urllib.parse.quote(topic)
        url = f"https://duckduckgo.com/html/?q={q}+tutorial"
        try:
            r = requests.get(url, timeout=10, headers={"User-Agent": "Mozilla/5.0"})
            r.raise_for_status()
            html_text = r.text
            items: List[Dict[str, str]] = []
            for m in re.finditer(r'<a[^>]*class="result__a"[^>]*href="([^"]+)"[^>]*>(.*?)</a>', html_text, flags=re.I|re.S):
                url = html.unescape(m.group(1))
                title = re.sub(r"<.*?>", "", html.unescape(m.group(2))).strip()
                if title and url and url.startswith("http"):
                    items.append({"title": title, "url": url})
                    if len(items) >= limit:
                        break
            return items
        except Exception:
            return []

    def recommend_youtube_ddg(self, topic: str, limit: int = 3) -> List[Dict[str, str]]:
        """Fetch YouTube video links from DuckDuckGo search."""
        q = urllib.parse.quote(f"{topic} site:youtube.com")
        url = f"https://duckduckgo.com/html/?q={q}"
        try:
            r = requests.get(url, timeout=10, headers={"User-Agent": "Mozilla/5.0"})
            r.raise_for_status()
            html_text = r.text
            items: List[Dict[str, str]] = []
            for m in re.finditer(r'<a[^>]*class="result__a"[^>]*href="([^"]+)"[^>]*>(.*?)</a>', html_text, flags=re.I|re.S):
                url = html.unescape(m.group(1))
                title = re.sub(r"<.*?>", "", html.unescape(m.group(2))).strip()
                if title and url and "youtube.com" in url and ("watch?v=" in url or "youtu.be/" in url):
                    items.append({"title": title, "url": url})
                    if len(items) >= limit:
                        break
            return items
        except Exception:
            return []

    def extract_topics_from_notes(self, notes_sections: List[str]) -> List[str]:
        """Extract topic headings (## or ###) from markdown sections."""
        topics: List[str] = []
        pattern = re.compile(r"^#{2,3}\s+(.+)$")
        for sec in notes_sections:
            for line in sec.splitlines():
                m = pattern.match(line.strip())
                if m:
                    t = m.group(1).strip()
                    if t and t not in topics:
                        topics.append(t)
        return topics

    def build_topic_context(self, notes_sections: List[str], topic: str) -> List[Dict[str, Any]]:
        """Build minimal context for LLM from notes related to a topic."""
        context_texts: List[str] = []
        for sec in notes_sections:
            if topic.lower() in sec.lower():
                context_texts.append(sec)
        if not context_texts:
            context_texts = notes_sections[:2]
        contexts = []
        for idx, t in enumerate(context_texts):
            contexts.append({"source": "notes", "chunk_index": idx, "text": t})
        return contexts

    def _generate_simple_notes(self, text: str, title: str):
        """Generate notes using the exact same implementation as LMS-AI-Assistant"""
        try:
            # Use the exact same note generation from LMS-AI-Assistant
            notes_sections = self.generate_notes_from_text(text, title=title)
            
            if not notes_sections:
                return f"# {title}\n\n*No content available to generate notes from.*"
            
            # Combine all sections
            return "\n\n".join(notes_sections)
            
        except Exception as e:
            return f"# {title}\n\n**Error generating notes:** {str(e)}\n\n```\n{text[:1000]}...\n```"
    
    def chunk_text(self, text: str, chunk_size: int = 800, chunk_overlap: int = 200) -> List[str]:
        """Split text into overlapping chunks suitable for embedding - exact copy from LMS-AI-Assistant"""
        if not text:
            return []

        normalized = " ".join(text.split())
        if len(normalized) <= chunk_size:
            return [normalized]

        chunks: List[str] = []
        start = 0
        while start < len(normalized):
            end = start + chunk_size
            chunk = normalized[start:end]
            chunks.append(chunk)
            if end >= len(normalized):
                break
            start = max(0, end - chunk_overlap)
        return chunks
    
    def generate_notes_from_text(self, text: str, title: str = None, chunk_size: int = 1200, chunk_overlap: int = 200) -> List[str]:
        """Generate markdown notes for the given text, chunk-by-chunk - exact copy from LMS-AI-Assistant"""
        chunks = self.chunk_text(text, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        if not chunks:
            return []

        backend = os.getenv("LLM_BACKEND")
        if backend:
            backend = backend.lower()
        else:
            backend = "groq" if os.getenv("GROQ_API_KEY") else "ollama"

        notes_sections: List[str] = []
        for ch in chunks:
            if backend == "ollama":
                md = self._llm_markdown_ollama(ch, title)
            else:
                md = self._llm_markdown_openai_compatible(ch, title)
            notes_sections.append(md or "")
        return notes_sections
    
    def iter_generate_notes_from_texts(self, texts: list[str], title: str = None, group_size: int = 3, pause_seconds: float = 2.0, max_retries: int = 3):
        """Yield markdown notes incrementally - exact copy from LMS-AI-Assistant"""
        backend = os.getenv("LLM_BACKEND")
        if backend:
            backend = backend.lower()
        else:
            backend = "groq" if os.getenv("GROQ_API_KEY") else "ollama"

        # 1) Deduplicate highly similar chunks with simple fingerprints
        def _fingerprint(t: str) -> str:
            import re, hashlib
            norm = re.sub(r"\W+", "", t.lower())
            return hashlib.sha1(norm.encode("utf-8")).hexdigest()

        seen: set[str] = set()
        unique_texts: list[str] = []
        for t in texts:
            fp = _fingerprint(t)
            if fp in seen:
                continue
            seen.add(fp)
            unique_texts.append(t)

        n = len(unique_texts)
        if n == 0:
            yield "No content available to generate notes from."
            return

        prev_outline: list[str] = []
        for i in range(0, n, max(1, group_size)):
            group = unique_texts[i : i + max(1, group_size)]
            content = "\n\n".join(group)
            
            # Skip if content is too short or empty
            if len(content.strip()) < 50:
                yield f"## Section {i//max(1, group_size) + 1}\n*Content too brief to process*"
                continue
                
            # Provide a brief outline of prior sections to discourage repetition
            if prev_outline:
                prefix = (
                    "Previously covered topics (do not repeat, only add new points):\n"
                    + "\n".join(f"- {o}" for o in prev_outline[-5:])
                    + "\n\n"
                )
            else:
                prefix = ""
            payload = prefix + content
            
            # Retry with simple backoff to avoid 429 limits
            attempt = 0
            md = ""
            while attempt <= max_retries:
                try:
                    if backend == "ollama":
                        md = self._llm_markdown_ollama(payload, title)
                    else:
                        md = self._llm_markdown_openai_compatible(payload, title)
                    break
                except Exception as e:
                    attempt += 1
                    if attempt > max_retries:
                        md = f"## Section {i//max(1, group_size) + 1}\n*Error generating notes: {str(e)[:100]}*"
                        break
                    time.sleep(min(pause_seconds * (1.5 ** (attempt - 1)), 15.0))
            
            # Ensure we always yield something
            if not md or len(md.strip()) < 10:
                md = f"## Section {i//max(1, group_size) + 1}\n*Generated content was empty or too short*"
            
            yield md
            
            # Extract a simple one-line summary as outline seed (first heading or first line)
            summary_line = (md.splitlines()[0] if md else "").strip()
            if summary_line:
                prev_outline.append(summary_line[:120])
            # Pause between sections to respect rate limits
            if pause_seconds > 0:
                time.sleep(pause_seconds)
    
    def _llm_markdown_ollama(self, chunk: str, title: str = None) -> str:
        """Generate markdown using Ollama - exact copy from LMS-AI-Assistant"""
        try:
            import ollama  # type: ignore
        except Exception as e:
            return f"## Error\n*Ollama not available: {str(e)}*"
        
        model_name = os.getenv("OLLAMA_MODEL", "llama3.1:8b")
        user_prompt = (
            (f"Title: {title}\n" if title else "")
            + "Create well-formatted lecture notes for the following content.\n\n"
            + chunk
        )
        messages = [
            {"role": "system", "content": self.NOTES_SYSTEM_PROMPT},
            {"role": "user", "content": user_prompt},
        ]
        try:
            resp = ollama.chat(model=model_name, messages=messages, options={"temperature": 0.2})
            return resp.get("message", {}).get("content", "")
        except Exception as e:
            return f"## Error\n*Ollama error: {str(e)}*"
    
    def _llm_markdown_openai_compatible(self, chunk: str, title: str = None) -> str:
        """Generate markdown using OpenAI-compatible API - exact copy from LMS-AI-Assistant"""
        try:
            from openai import OpenAI
        except Exception as e:
            return f"## Error\n*OpenAI not available: {str(e)}*"

        api_key = os.getenv("GROQ_API_KEY") or os.getenv("OPENAI_API_KEY")
        if not api_key:
            return f"## Error\n*Missing GROQ_API_KEY/OPENAI_API_KEY for notes generation*"
        
        base_url = os.getenv("GROQ_BASE_URL", "https://api.groq.com/openai/v1")
        client = OpenAI(api_key=api_key, base_url=base_url)

        model_name = os.getenv("GROQ_MODEL")
        if not model_name:
            try:
                models = client.models.list()
                llama_models = [m.id for m in getattr(models, "data", []) if "llama" in getattr(m, "id", "")]
                model_name = sorted(llama_models)[0] if llama_models else "llama-3.1-8b-instant"
            except Exception:
                model_name = "llama-3.1-8b-instant"

        user_prompt = (
            (f"Title: {title}\n" if title else "")
            + "Create well-formatted lecture notes for the following content.\n\n"
            + chunk
        )
        messages = [
            {"role": "system", "content": self.NOTES_SYSTEM_PROMPT},
            {"role": "user", "content": user_prompt},
        ]
        try:
            resp = client.chat.completions.create(
                model=model_name,
                messages=messages,
                temperature=0.2,
                max_tokens=800,
            )
            return resp.choices[0].message.content or ""
        except Exception as e:
            return f"## Error\n*API error: {str(e)}*"
    
    # Notes system prompt - exact copy from LMS-AI-Assistant
    NOTES_SYSTEM_PROMPT = (
        "You are an expert technical writer creating lecture-article notes for students. "
        "Only include content a diligent student would write in a notebook: crisp definitions, key formulas, step-by-step procedures, concise examples, short summaries, caveats, and essential code snippets. "
        "Exclude fluff, marketing, anecdotes, and repeated text.\n\n"
        "Formatting rules:\n"
        "- Use clear H2/H3 headings (##, ###)\n"
        "- Prefer short bullet lists\n"
        "- Put code in fenced blocks with language hints when apparent (```python, ```js, etc.)\n"
        "- Keep paragraphs short and focused\n"
        "- Do not hallucinate; only use the provided content.\n"
    )
    
    def _is_section_header(self, line: str, index: int, all_lines: List[str]) -> bool:
        """Detect if a line is likely a section header"""
        # Common header patterns
        header_patterns = [
            r'^#{1,6}\s+',  # Markdown headers
            r'^\d+\.?\s+[A-Z]',  # Numbered sections
            r'^[A-Z][A-Z\s]{3,}$',  # ALL CAPS headers
            r'^[A-Z][a-z]+\s+[A-Z]',  # Title Case headers
            r'^\d+\.\d+',  # Decimal numbering
            r'^Chapter\s+\d+',  # Chapter headers
            r'^Section\s+\d+',  # Section headers
            r'^Part\s+\d+',  # Part headers
        ]
        
        import re
        for pattern in header_patterns:
            if re.match(pattern, line):
                return True
        
        # Check if line is significantly shorter than average and followed by content
        if len(line) < 50 and index < len(all_lines) - 1:
            next_line = all_lines[index + 1].strip()
            if next_line and len(next_line) > len(line) * 1.5:
                return True
        
        return False
    
    def _clean_header(self, line: str) -> str:
        """Clean and format section headers"""
        import re
        
        # Remove markdown headers
        line = re.sub(r'^#{1,6}\s+', '', line)
        
        # Remove numbering
        line = re.sub(r'^\d+\.?\s*', '', line)
        
        # Clean up formatting
        line = line.strip('*_`')
        
        # Convert to title case if all caps
        if line.isupper() and len(line) > 3:
            line = line.title()
        
        return line
    
    def _format_section(self, title: str, content: List[str]) -> str:
        """Format a section with proper markdown structure"""
        section = f"## 📖 {title}\n\n"
        
        # Process content within section
        processed_content = []
        in_list = False
        in_code = False
        
        for line in content:
            if not line.strip():
                if in_list:
                    processed_content.append("")
                continue
            
            # Detect code blocks
            if line.strip().startswith('```'):
                in_code = not in_code
                processed_content.append(line)
                continue
            
            if in_code:
                processed_content.append(line)
                continue
            
            # Detect lists
            if self._is_list_item(line):
                if not in_list:
                    processed_content.append("")
                processed_content.append(f"- {self._clean_list_item(line)}")
                in_list = True
            else:
                if in_list:
                    processed_content.append("")
                in_list = False
                
                # Format different types of content
                formatted_line = self._format_content_line(line)
                processed_content.append(formatted_line)
        
        section += '\n'.join(processed_content)
        return section
    
    def _is_list_item(self, line: str) -> bool:
        """Check if line is a list item"""
        import re
        list_patterns = [
            r'^\s*[-*+]\s+',  # Bullet points
            r'^\s*\d+\.\s+',  # Numbered lists
            r'^\s*[a-zA-Z]\.\s+',  # Letter lists
        ]
        
        for pattern in list_patterns:
            if re.match(pattern, line):
                return True
        return False
    
    def _clean_list_item(self, line: str) -> str:
        """Clean list item formatting"""
        import re
        # Remove bullet points and numbering
        line = re.sub(r'^\s*[-*+]\s+', '', line)
        line = re.sub(r'^\s*\d+\.\s+', '', line)
        line = re.sub(r'^\s*[a-zA-Z]\.\s+', '', line)
        return line.strip()
    
    def _format_content_line(self, line: str) -> str:
        """Format individual content lines with appropriate styling"""
        line_lower = line.lower()
        
        # Definitions
        if any(keyword in line_lower for keyword in ['definition', 'define', 'what is', 'means']):
            return f"**📖 Definition:** {line}"
        
        # Examples
        if any(keyword in line_lower for keyword in ['example', 'for instance', 'such as', 'e.g.']):
            return f"**💡 Example:** {line}"
        
        # Important notes
        if any(keyword in line_lower for keyword in ['important', 'note', 'remember', 'key', 'critical']):
            return f"⚠️ **Important:** {line}"
        
        # Steps or procedures
        if any(keyword in line_lower for keyword in ['step', 'process', 'procedure', 'method']):
            return f"**🔧 Step:** {line}"
        
        # Formulas or equations
        if any(char in line for char in ['=', '+', '-', '*', '/', '^']) and len(line) < 100:
            return f"**🧮 Formula:** `{line}`"
        
        # Questions
        if line.strip().endswith('?'):
            return f"**❓ Question:** {line}"
        
        # Regular content
        return line
    
    def _format_general_content(self, title: str, lines: List[str]) -> str:
        """Format content when no clear sections are detected"""
        section = f"## 📄 Content Overview\n\n"
        
        # Group related lines
        current_group = []
        groups = []
        
        for line in lines:
            if not line.strip():
                if current_group:
                    groups.append(current_group)
                    current_group = []
            else:
                current_group.append(line)
        
        if current_group:
            groups.append(current_group)
        
        # Format each group
        for i, group in enumerate(groups):
            if len(group) == 1:
                section += f"{self._format_content_line(group[0])}\n\n"
            else:
                section += f"**Topic {i+1}:**\n"
                for line in group:
                    section += f"- {self._format_content_line(line)}\n"
                section += "\n"
        
        return section
    
    def _generate_summary(self, text: str) -> str:
        """Generate a brief summary of the content"""
        # Simple extractive summary (first few sentences)
        sentences = text.split('. ')
        if len(sentences) <= 3:
            return text
        else:
            summary_sentences = sentences[:3]
            return '. '.join(summary_sentences) + "..."
    
    def _generate_metadata(self, title: str, text: str) -> str:
        """Generate document metadata"""
        word_count = len(text.split())
        char_count = len(text)
        line_count = len(text.split('\n'))
        
        return f"""
- **Title:** {title}
- **Word Count:** {word_count:,}
- **Character Count:** {char_count:,}
- **Lines:** {line_count:,}
- **Estimated Reading Time:** {word_count // 200 + 1} minutes
"""
    
    def _format_content_preview(self, text: str, filename: str) -> str:
        """Format content preview with enhanced styling"""
        try:
            # Basic text processing for preview
            lines = text.split('\n')
            processed_lines = []
            
            # Add document header
            processed_lines.append(f"# 📄 {filename}")
            processed_lines.append("")
            processed_lines.append("*Content Preview*")
            processed_lines.append("")
            
            # Process content with basic formatting
            for line in lines[:50]:  # Limit to first 50 lines for preview
                line = line.strip()
                if not line:
                    processed_lines.append("")
                    continue
                
                # Apply basic formatting
                formatted_line = self._format_content_line(line)
                processed_lines.append(formatted_line)
            
            # Add truncation notice if content is long
            if len(lines) > 50:
                processed_lines.append("")
                processed_lines.append("---")
                processed_lines.append("*Content truncated for preview...*")
            
            return '\n'.join(processed_lines)
            
        except Exception as e:
            return f"# 📄 {filename}\n\n**Error formatting preview:** {str(e)}\n\n```\n{text[:500]}...\n```"

    def load_course_data(self, course_id: int):
        """Load comprehensive course data"""
        if not st.session_state.client:
            return {}
        
        try:
            with st.spinner("Loading course data..."):
                course_data = {
                    'assignments': [],
                    'modules': [],
                    'users': [],
                    'announcements': []
                }
                
                # Load assignments
                try:
                    assignments = st.session_state.client.get_course_assignments(course_id)
                    course_data['assignments'] = assignments
                except Exception as e:
                    logger.warning(f"Could not load assignments: {e}")
                
                # Load modules
                try:
                    modules = st.session_state.client.get_course_modules(course_id, include_items=True)
                    course_data['modules'] = modules
                except Exception as e:
                    logger.warning(f"Could not load modules: {e}")
                
                # Load course files from local downloads directory instead of Canvas API
                course_data['files'] = self._load_local_course_files(course_id)
                
                # Load users
                try:
                    users = st.session_state.client.get_course_users(course_id)
                    course_data['users'] = users
                except Exception as e:
                    logger.warning(f"Could not load users: {e}")
                
                # Load announcements
                try:
                    announcements = []
                    page = 1
                    per_page = 50
                    
                    while True:
                        params = {
                            'per_page': per_page,
                            'page': page,
                            'context_codes[]': f'course_{course_id}'
                        }
                        
                        response = st.session_state.client._make_request(
                            'GET', 
                            '/api/v1/announcements', 
                            params=params
                        )
                        
                        if not response:
                            break
                        
                        announcements.extend(response)
                        
                        if len(response) < per_page:
                            break
                        
                        page += 1
                    
                    course_data['announcements'] = announcements
                except Exception as e:
                    logger.warning(f"Could not load announcements: {e}")
                
                return course_data
                
        except Exception as e:
            st.error(f"❌ Failed to load course data: {str(e)}")
            logger.error(f"Load course data error: {e}")
            return {}
    
    def render_course_selection(self):
        """Render course selection interface"""
        if not st.session_state.authenticated:
            return
        
        st.header("📚 Course Selection")
        
        # Load courses button
        col1, col2 = st.columns([1, 3])
        with col1:
            if st.button("🔄 Refresh Courses", type="primary"):
                self.load_courses()
        
        with col2:
            if st.session_state.courses:
                st.success(f"Found {len(st.session_state.courses)} courses")
        
        # Course selection
        if st.session_state.courses:
            course_options = {
                f"{course.get('name', 'Unknown')} (ID: {course.get('id')})": course
                for course in st.session_state.courses
            }
            
            selected_course_name = st.selectbox(
                "Select a course:",
                options=list(course_options.keys()),
                index=0
            )
            
            if selected_course_name:
                st.session_state.selected_course = course_options[selected_course_name]
                # Load course data when course is selected
                if st.session_state.selected_course:
                    course_id = st.session_state.selected_course.get('id')
                    st.session_state.course_data = self.load_course_data(course_id)
        else:
            st.info("Click 'Refresh Courses' to load your courses")
    
    def render_workflow_status(self):
        """Render enhanced workflow progress indicator with actionable insights"""
        global NAVDEEP_AVAILABLE
        course = st.session_state.selected_course
        course_data = st.session_state.course_data
        
        # Calculate comprehensive workflow metrics
        download_path = f"/Users/praveenbhandari/sf-vibe/downloads/{course.get('name', '').replace(' ', '_').replace('/', '_')}"
        files_exist = os.path.exists(download_path) and len(os.listdir(download_path)) > 0
        
        # Count downloadable content
        total_files = len(course_data.get('files', []))
        total_assignments = len(course_data.get('assignments', []))
        total_modules = len(course_data.get('modules', []))
        
        # Check permission status
        files_accessible = total_files > 0 or not any('forbidden' in str(e).lower() or 'insufficient permissions' in str(e).lower() 
                                                    for e in [getattr(course_data, 'files_error', None)] if e)
        
        # Calculate progress percentage
        progress_steps = [True, True, files_accessible, False]  # login, course, files, ai
        progress_percent = (sum(progress_steps) / len(progress_steps)) * 100
        
        # Enhanced workflow container
        st.markdown(f"""
        <div style="
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            padding: 2rem;
            border-radius: 15px;
            margin-bottom: 2rem;
            box-shadow: 0 8px 32px rgba(0,0,0,0.1);
            color: white;
        ">
            <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 1.5rem;">
                <h3 style="margin: 0; color: white; font-size: 1.5rem;">🚀 Workflow Dashboard</h3>
                <div style="
                    background: rgba(255,255,255,0.2);
                    padding: 0.5rem 1rem;
                    border-radius: 20px;
                    font-weight: bold;
                ">{progress_percent:.0f}% Complete</div>
            </div>
            
            <!-- Progress Bar -->
            <div style="
                background: rgba(255,255,255,0.2);
                height: 8px;
                border-radius: 4px;
                margin-bottom: 1rem;
                overflow: hidden;
            ">
                <div style="
                    background: linear-gradient(90deg, #00d4aa, #00b894);
                    height: 100%;
                    width: {progress_percent}%;
                    border-radius: 4px;
                    transition: width 0.3s ease;
                "></div>
            </div>
            
            <!-- Permission Status Indicator -->
            <div style="
                background: rgba(255,255,255,0.1);
                padding: 10px;
                border-radius: 8px;
                margin-bottom: 1rem;
                border-left: 3px solid {'#4CAF50' if files_accessible else '#FF9800'};
            ">
                <p style="color: rgba(255,255,255,0.9); margin: 0; font-size: 14px; font-weight: 500;">
                    {'✅ Files Access: Available' if files_accessible else '⚠️ Files Access: Restricted - Check permissions'}
                </p>
                {f'<p style="color: rgba(255,255,255,0.7); margin: 5px 0 0 0; font-size: 12px;">Found {total_files} files</p>' if files_accessible else ''}
            </div>
        </div>
        """, unsafe_allow_html=True)
        
        # Enhanced progress steps with metrics
        steps_data = [
            {
                "icon": "✅",
                "title": "Authentication",
                "status": "Connected",
                "color": "#00d4aa",
                "completed": True,
                "metric": f"User: {st.session_state.user_info.get('name', 'Unknown')[:15]}...",
                "action": None
            },
            {
                "icon": "✅",
                "title": "Course Selection",
                "status": "Active",
                "color": "#00d4aa",
                "completed": True,
                "metric": f"{total_assignments} assignments, {total_modules} modules",
                "action": None
            },
            {
                "icon": "✅" if files_exist else "⏳",
                "title": "Content Download",
                "status": "Ready" if files_exist else "Pending",
                "color": "#00d4aa" if files_exist else "#ffd93d",
                "completed": files_exist,
                "metric": f"{total_files} files available",
                "action": "download" if not files_exist else None
            },
            {
                "icon": "🤖" if NAVDEEP_AVAILABLE else "❌",
                "title": "AI Analysis",
                "status": "Ready" if NAVDEEP_AVAILABLE and files_exist else "Waiting",
                "color": "#74b9ff" if NAVDEEP_AVAILABLE and files_exist else "#636e72",
                "completed": False,
                "metric": "AI components available" if NAVDEEP_AVAILABLE else "AI unavailable",
                "action": "ai_notes" if NAVDEEP_AVAILABLE and files_exist else None
            }
        ]
        
        # Render enhanced step cards
        cols = st.columns(4)
        for i, step in enumerate(steps_data):
            with cols[i]:
                # Determine card style based on completion
                card_bg = "rgba(0, 212, 170, 0.1)" if step["completed"] else "rgba(255, 255, 255, 0.05)"
                border_color = step["color"] if step["completed"] else "#e0e0e0"
                
                st.markdown(f"""
                <div style="
                    background: {card_bg};
                    border: 2px solid {border_color};
                    border-radius: 12px;
                    padding: 1.5rem;
                    text-align: center;
                    height: 180px;
                    display: flex;
                    flex-direction: column;
                    justify-content: space-between;
                    transition: all 0.3s ease;
                    margin-bottom: 1rem;
                ">
                    <div>
                        <div style="font-size: 2.5rem; margin-bottom: 0.5rem;">{step["icon"]}</div>
                        <div style="font-weight: 700; color: {step["color"]}; margin-bottom: 0.3rem; font-size: 0.9rem;">{step["title"]}</div>
                        <div style="font-size: 0.75rem; color: #666; margin-bottom: 0.5rem;">{step["status"]}</div>
                    </div>
                    <div style="font-size: 0.7rem; color: #888; line-height: 1.2;">{step["metric"]}</div>
                </div>
                """, unsafe_allow_html=True)
                
                # Add action buttons for pending steps
                if step["action"] == "download":
                    if st.button("📥 Start Download", key=f"action_download_{i}", type="primary", use_container_width=True):
                        st.switch_page("pages/3_📥_Download_Files.py")
                elif step["action"] == "ai_notes":
                    if st.button("🤖 Generate Notes", key=f"action_ai_{i}", type="primary", use_container_width=True):
                        st.switch_page("pages/4_🤖_AI_Notes.py")
        
        # Next steps recommendations
        st.markdown("### 🎯 Recommended Next Steps")
        
        if not files_exist:
            st.info("💡 **Download course files** to unlock AI-powered note generation and content analysis.")
        elif NAVDEEP_AVAILABLE and files_exist:
            st.success("🚀 **Ready for AI analysis!** Generate intelligent notes from your downloaded content.")
        elif not NAVDEEP_AVAILABLE:
            st.warning("⚠️ **AI components unavailable.** Some advanced features may be limited.")
        else:
            st.success("✨ **All systems ready!** Explore your course content and generate insights.")
    
    def render_course_overview(self):
        """Render enhanced course overview with comprehensive metrics and insights"""
        if not st.session_state.selected_course:
            return
        
        course = st.session_state.selected_course
        course_data = st.session_state.course_data
        
        # Calculate comprehensive course metrics
        assignments = course_data.get('assignments', [])
        modules = course_data.get('modules', [])
        files = course_data.get('files', [])
        users = course_data.get('users', [])
        announcements = course_data.get('announcements', [])
        
        # Advanced analytics
        total_points = sum(a.get('points_possible', 0) for a in assignments if a.get('points_possible'))
        due_assignments = [a for a in assignments if a.get('due_at')]
        overdue_count = 0
        upcoming_count = 0
        
        # Calculate due date statistics
        from datetime import datetime, timezone
        now = datetime.now(timezone.utc)
        
        for assignment in due_assignments:
            try:
                due_date = datetime.fromisoformat(assignment['due_at'].replace('Z', '+00:00'))
                if due_date < now:
                    overdue_count += 1
                elif (due_date - now).days <= 7:
                    upcoming_count += 1
            except:
                continue
        
        # Calculate file size statistics
        total_file_size = sum(f.get('size', 0) for f in files)
        file_size_mb = total_file_size / (1024 * 1024) if total_file_size > 0 else 0
        
        # Course header with enhanced styling
        st.markdown(f"""
        <div style="
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            padding: 2rem;
            border-radius: 15px;
            margin-bottom: 2rem;
            color: white;
            box-shadow: 0 8px 32px rgba(0,0,0,0.1);
        ">
            <h2 style="margin: 0 0 0.5rem 0; color: white; font-size: 2rem;">📖 {course.get('name', 'Unknown Course')}</h2>
            <p style="margin: 0; opacity: 0.9; font-size: 1.1rem;">{course.get('course_code', 'N/A')} • {course.get('term', {}).get('name', 'N/A')}</p>
        </div>
        """, unsafe_allow_html=True)
        
        # Enhanced metrics dashboard
        st.markdown("### 📊 Course Analytics Dashboard")
        
        # Primary metrics row
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #ff6b6b, #ee5a52);
                padding: 1.5rem;
                border-radius: 12px;
                text-align: center;
                color: white;
                box-shadow: 0 4px 15px rgba(255,107,107,0.3);
            ">
                <div style="font-size: 2.5rem; margin-bottom: 0.5rem;">📝</div>
                <div style="font-size: 2rem; font-weight: bold; margin-bottom: 0.25rem;">{}</div>
                <div style="font-size: 0.9rem; opacity: 0.9;">Assignments</div>
                <div style="font-size: 0.8rem; margin-top: 0.5rem;">{} points total</div>
            </div>
            """.format(len(assignments), int(total_points)), unsafe_allow_html=True)
        
        with col2:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #4ecdc4, #44a08d);
                padding: 1.5rem;
                border-radius: 12px;
                text-align: center;
                color: white;
                box-shadow: 0 4px 15px rgba(78,205,196,0.3);
            ">
                <div style="font-size: 2.5rem; margin-bottom: 0.5rem;">📚</div>
                <div style="font-size: 2rem; font-weight: bold; margin-bottom: 0.25rem;">{}</div>
                <div style="font-size: 0.9rem; opacity: 0.9;">Modules</div>
                <div style="font-size: 0.8rem; margin-top: 0.5rem;">{} total items</div>
            </div>
            """.format(len(modules), sum(len(m.get('items', [])) for m in modules)), unsafe_allow_html=True)
        
        with col3:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #45b7d1, #3498db);
                padding: 1.5rem;
                border-radius: 12px;
                text-align: center;
                color: white;
                box-shadow: 0 4px 15px rgba(69,183,209,0.3);
            ">
                <div style="font-size: 2.5rem; margin-bottom: 0.5rem;">📁</div>
                <div style="font-size: 2rem; font-weight: bold; margin-bottom: 0.25rem;">{}</div>
                <div style="font-size: 0.9rem; opacity: 0.9;">Files</div>
                <div style="font-size: 0.8rem; margin-top: 0.5rem;">{:.1f} MB total</div>
            </div>
            """.format(len(files), file_size_mb), unsafe_allow_html=True)
        
        with col4:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #a8e6cf, #7fcdcd);
                padding: 1.5rem;
                border-radius: 12px;
                text-align: center;
                color: white;
                box-shadow: 0 4px 15px rgba(168,230,207,0.3);
            ">
                <div style="font-size: 2.5rem; margin-bottom: 0.5rem;">👥</div>
                <div style="font-size: 2rem; font-weight: bold; margin-bottom: 0.25rem;">{}</div>
                <div style="font-size: 0.9rem; opacity: 0.9;">Students</div>
                <div style="font-size: 0.8rem; margin-top: 0.5rem;">{} announcements</div>
            </div>
            """.format(len(users), len(announcements)), unsafe_allow_html=True)
        
        st.markdown("<br>", unsafe_allow_html=True)
        
        # Assignment insights row
        if assignments:
            st.markdown("### ⏰ Assignment Insights")
            col1, col2, col3 = st.columns(3)
            
            with col1:
                st.markdown("""
                <div style="
                    background: linear-gradient(135deg, #ffd93d, #ffb74d);
                    padding: 1.5rem;
                    border-radius: 12px;
                    text-align: center;
                    color: white;
                    box-shadow: 0 4px 15px rgba(255,217,61,0.3);
                ">
                    <div style="font-size: 2rem; margin-bottom: 0.5rem;">⚠️</div>
                    <div style="font-size: 1.8rem; font-weight: bold; margin-bottom: 0.25rem;">{}</div>
                    <div style="font-size: 0.9rem; opacity: 0.9;">Overdue</div>
                </div>
                """.format(overdue_count), unsafe_allow_html=True)
            
            with col2:
                st.markdown("""
                <div style="
                    background: linear-gradient(135deg, #74b9ff, #0984e3);
                    padding: 1.5rem;
                    border-radius: 12px;
                    text-align: center;
                    color: white;
                    box-shadow: 0 4px 15px rgba(116,185,255,0.3);
                ">
                    <div style="font-size: 2rem; margin-bottom: 0.5rem;">📅</div>
                    <div style="font-size: 1.8rem; font-weight: bold; margin-bottom: 0.25rem;">{}</div>
                    <div style="font-size: 0.9rem; opacity: 0.9;">Due This Week</div>
                </div>
                """.format(upcoming_count), unsafe_allow_html=True)
            
            with col3:
                completion_rate = ((len(assignments) - overdue_count) / len(assignments) * 100) if assignments else 0
                st.markdown("""
                <div style="
                    background: linear-gradient(135deg, #00b894, #00a085);
                    padding: 1.5rem;
                    border-radius: 12px;
                    text-align: center;
                    color: white;
                    box-shadow: 0 4px 15px rgba(0,184,148,0.3);
                ">
                    <div style="font-size: 2rem; margin-bottom: 0.5rem;">✅</div>
                    <div style="font-size: 1.8rem; font-weight: bold; margin-bottom: 0.25rem;">{:.0f}%</div>
                    <div style="font-size: 0.9rem; opacity: 0.9;">On Track</div>
                </div>
                """.format(completion_rate), unsafe_allow_html=True)
        
        st.markdown("<br>", unsafe_allow_html=True)
        
        # Quick actions section
        st.markdown("### 🚀 Quick Actions")
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            if st.button("📝 View Assignments", type="primary", use_container_width=True):
                st.session_state.active_tab = "assignments"
                st.rerun()
        
        with col2:
            if st.button("📚 Browse Modules", type="secondary", use_container_width=True):
                st.session_state.active_tab = "modules"
                st.rerun()
        
        with col3:
            if st.button("📁 Download Files", type="secondary", use_container_width=True):
                st.session_state.active_tab = "download"
                st.rerun()
        
        with col4:
            if st.button("📊 Export Data", type="secondary", use_container_width=True):
                st.session_state.active_tab = "export"
                st.rerun()
        
        # Detailed course information in expandable section
        with st.expander("📋 Detailed Course Information", expanded=False):
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown("**📚 Course Details**")
                st.write(f"• **Course ID:** {course.get('id')}")
                st.write(f"• **Course Code:** {course.get('course_code', 'N/A')}")
                st.write(f"• **Term:** {course.get('term', {}).get('name', 'N/A')}")
                st.write(f"• **Workflow State:** {course.get('workflow_state', 'N/A')}")
                
            with col2:
                st.markdown("**📅 Timeline**")
                st.write(f"• **Start Date:** {course.get('start_at', 'N/A')[:10] if course.get('start_at') else 'N/A'}")
                st.write(f"• **End Date:** {course.get('end_at', 'N/A')[:10] if course.get('end_at') else 'N/A'}")
                st.write(f"• **Created:** {course.get('created_at', 'N/A')[:10] if course.get('created_at') else 'N/A'}")
                st.write(f"• **Updated:** {course.get('updated_at', 'N/A')[:10] if course.get('updated_at') else 'N/A'}")
            
            if course.get('description'):
                st.markdown("**📝 Course Description**")
                st.write(course.get('description'))
    
    def render_assignments_section(self):
        """Render enhanced assignments section with improved visual design"""
        if not st.session_state.course_data.get('assignments'):
            return
        
        # Enhanced header with gradient background
        st.markdown("""
        <div style="
            background: linear-gradient(135deg, #ff6b6b 0%, #ee5a52 100%);
            padding: 1.5rem;
            border-radius: 12px;
            margin-bottom: 1.5rem;
            color: white;
            box-shadow: 0 4px 15px rgba(255,107,107,0.3);
        ">
            <h2 style="margin: 0; color: white; font-size: 1.8rem;">📝 Course Assignments</h2>
            <p style="margin: 0.5rem 0 0 0; opacity: 0.9;">Manage and track your assignment progress</p>
        </div>
        """, unsafe_allow_html=True)
        
        assignments = st.session_state.course_data['assignments']
        
        # Enhanced filter section with better styling
        st.markdown("### 🔧 Filter & Sort Options")
        col1, col2, col3 = st.columns([1, 1, 1])
        
        with col1:
            show_completed = st.checkbox("✅ Show completed assignments", value=True)
        
        with col2:
            sort_by = st.selectbox("📊 Sort by", ["Due Date", "Name", "Points"])
        
        with col3:
            # Add assignment status filter
            status_filter = st.selectbox("🎯 Status Filter", ["All", "Upcoming", "Overdue", "No Due Date"])
        
        # Calculate assignment statistics for better insights
        from datetime import datetime, timezone
        now = datetime.now(timezone.utc)
        
        overdue_assignments = []
        upcoming_assignments = []
        no_due_date = []
        
        for assignment in assignments:
            due_date = assignment.get('due_at')
            if not due_date:
                no_due_date.append(assignment)
            else:
                try:
                    dt = datetime.fromisoformat(due_date.replace('Z', '+00:00'))
                    if dt < now:
                        overdue_assignments.append(assignment)
                    elif (dt - now).days <= 7:
                        upcoming_assignments.append(assignment)
                except:
                    no_due_date.append(assignment)
        
        # Assignment statistics dashboard
        st.markdown("### 📊 Assignment Overview")
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #74b9ff, #0984e3);
                padding: 1rem;
                border-radius: 10px;
                text-align: center;
                color: white;
                box-shadow: 0 3px 10px rgba(116,185,255,0.3);
            ">
                <div style="font-size: 1.8rem; margin-bottom: 0.3rem;">📝</div>
                <div style="font-size: 1.5rem; font-weight: bold;">{}</div>
                <div style="font-size: 0.8rem; opacity: 0.9;">Total</div>
            </div>
            """.format(len(assignments)), unsafe_allow_html=True)
        
        with col2:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #ffd93d, #ffb74d);
                padding: 1rem;
                border-radius: 10px;
                text-align: center;
                color: white;
                box-shadow: 0 3px 10px rgba(255,217,61,0.3);
            ">
                <div style="font-size: 1.8rem; margin-bottom: 0.3rem;">⚠️</div>
                <div style="font-size: 1.5rem; font-weight: bold;">{}</div>
                <div style="font-size: 0.8rem; opacity: 0.9;">Overdue</div>
            </div>
            """.format(len(overdue_assignments)), unsafe_allow_html=True)
        
        with col3:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #00b894, #00a085);
                padding: 1rem;
                border-radius: 10px;
                text-align: center;
                color: white;
                box-shadow: 0 3px 10px rgba(0,184,148,0.3);
            ">
                <div style="font-size: 1.8rem; margin-bottom: 0.3rem;">📅</div>
                <div style="font-size: 1.5rem; font-weight: bold;">{}</div>
                <div style="font-size: 0.8rem; opacity: 0.9;">This Week</div>
            </div>
            """.format(len(upcoming_assignments)), unsafe_allow_html=True)
        
        with col4:
            total_points = sum(a.get('points_possible', 0) for a in assignments if a.get('points_possible'))
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #a8e6cf, #7fcdcd);
                padding: 1rem;
                border-radius: 10px;
                text-align: center;
                color: white;
                box-shadow: 0 3px 10px rgba(168,230,207,0.3);
            ">
                <div style="font-size: 1.8rem; margin-bottom: 0.3rem;">🎯</div>
                <div style="font-size: 1.5rem; font-weight: bold;">{}</div>
                <div style="font-size: 0.8rem; opacity: 0.9;">Total Points</div>
            </div>
            """.format(int(total_points)), unsafe_allow_html=True)
        
        st.markdown("<br>", unsafe_allow_html=True)
        
        # Filter and sort assignments
        filtered_assignments = assignments.copy()
        
        if not show_completed:
            filtered_assignments = [a for a in filtered_assignments if not a.get('has_submitted_submissions', False)]
        
        # Apply status filter
        if status_filter == "Overdue":
            filtered_assignments = [a for a in filtered_assignments if a in overdue_assignments]
        elif status_filter == "Upcoming":
            filtered_assignments = [a for a in filtered_assignments if a in upcoming_assignments]
        elif status_filter == "No Due Date":
            filtered_assignments = [a for a in filtered_assignments if a in no_due_date]
        
        if sort_by == "Due Date":
            filtered_assignments.sort(key=lambda x: x.get('due_at') or '9999-12-31')
        elif sort_by == "Name":
            filtered_assignments.sort(key=lambda x: x.get('name', ''))
        elif sort_by == "Points":
            filtered_assignments.sort(key=lambda x: x.get('points_possible') or 0, reverse=True)
        
        # Display assignments with enhanced card design
        if not filtered_assignments:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #ddd6fe, #c4b5fd);
                padding: 2rem;
                border-radius: 12px;
                text-align: center;
                color: #6b46c1;
                margin: 1rem 0;
            ">
                <div style="font-size: 3rem; margin-bottom: 1rem;">📝</div>
                <h3 style="margin: 0; color: #6b46c1;">No assignments found</h3>
                <p style="margin: 0.5rem 0 0 0; opacity: 0.8;">Try adjusting your filter criteria</p>
            </div>
            """, unsafe_allow_html=True)
        else:
            st.markdown(f"### 📋 Assignment List ({len(filtered_assignments)} found)")
            
            for i, assignment in enumerate(filtered_assignments, 1):
                # Determine card color based on assignment status
                due_date = assignment.get('due_at')
                card_color = "#f8f9fa"  # default
                border_color = "#4CAF50"  # default green
                status_badge = "📝"
                status_text = "Active"
                
                if due_date:
                    try:
                        dt = datetime.fromisoformat(due_date.replace('Z', '+00:00'))
                        if dt < now:
                            border_color = "#ff6b6b"  # red for overdue
                            status_badge = "⚠️"
                            status_text = "Overdue"
                        elif (dt - now).days <= 3:
                            border_color = "#ffd93d"  # yellow for due soon
                            status_badge = "⏰"
                            status_text = "Due Soon"
                    except:
                        pass
                
                with st.container():
                    st.markdown(
                        f"""
                        <div style="
                            background: linear-gradient(135deg, {card_color}, #ffffff);
                            border-left: 6px solid {border_color};
                            padding: 1.5rem;
                            margin: 1rem 0;
                            border-radius: 0 12px 12px 0;
                            box-shadow: 0 4px 15px rgba(0,0,0,0.1);
                            transition: transform 0.2s ease;
                        ">
                            <div style="
                                display: flex;
                                justify-content: space-between;
                                align-items: center;
                                margin-bottom: 1rem;
                            ">
                                <h3 style="margin: 0; color: #2c3e50; font-size: 1.3rem;">{status_badge} {assignment.get('name', 'No name')}</h3>
                                <span style="
                                    background: {border_color};
                                    color: white;
                                    padding: 0.3rem 0.8rem;
                                    border-radius: 20px;
                                    font-size: 0.8rem;
                                    font-weight: bold;
                                ">{status_text}</span>
                            </div>
                        </div>
                        """, 
                        unsafe_allow_html=True
                    )
                    
                    col1, col2, col3 = st.columns([2, 1, 1])
                    
                    with col1:
                        # Description with better formatting
                        if assignment.get('description'):
                            desc = assignment.get('description', '')
                            if len(desc) > 150:
                                with st.expander("📄 View Full Description"):
                                    st.markdown(desc, unsafe_allow_html=True)
                            else:
                                st.markdown(f"*{desc[:150]}...*" if len(desc) > 150 else f"*{desc}*")
                        else:
                            st.markdown("*No description available*")
                    
                    with col2:
                        if due_date:
                            try:
                                dt = datetime.fromisoformat(due_date.replace('Z', '+00:00'))
                                formatted_date = dt.strftime('%b %d, %Y')
                                formatted_time = dt.strftime('%I:%M %p')
                                
                                # Calculate days until due
                                days_diff = (dt - now).days
                                if days_diff < 0:
                                    time_text = f"⚠️ {abs(days_diff)} days overdue"
                                elif days_diff == 0:
                                    time_text = "🔥 Due today!"
                                elif days_diff <= 7:
                                    time_text = f"⏰ Due in {days_diff} days"
                                else:
                                    time_text = f"📅 Due in {days_diff} days"
                                
                                st.markdown(f"**📅 Due Date**")
                                st.markdown(f"{formatted_date}")
                                st.markdown(f"{formatted_time}")
                                st.markdown(f"*{time_text}*")
                            except:
                                st.markdown(f"**📅 Due Date**")
                                st.markdown(f"{due_date[:10]}")
                        else:
                            st.markdown(f"**📅 Due Date**")
                            st.markdown("⏰ No due date set")
                    
                    with col3:
                        points = assignment.get('points_possible')
                        st.markdown(f"**🎯 Points**")
                        if points is not None and points > 0:
                            st.markdown(f"⭐ **{points}** pts")
                            # Add points visualization
                            if points >= 100:
                                st.markdown("🏆 *High value*")
                            elif points >= 50:
                                st.markdown("🥈 *Medium value*")
                            else:
                                st.markdown("🥉 *Low value*")
                        else:
                            st.markdown("➖ Not graded")
                    
                    # Assignment actions
                    if assignment.get('html_url'):
                        st.markdown(f"[🔗 Open Assignment]({assignment.get('html_url')})")
                    
                    st.markdown("<br>", unsafe_allow_html=True)
    
    def render_modules_section(self):
        """Render enhanced modules section with improved visual design and color coding"""
        if not st.session_state.course_data.get('modules'):
            return
        
        # Enhanced header with gradient background
        st.markdown("""
        <div style="
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            padding: 20px;
            border-radius: 10px;
            margin-bottom: 20px;
            box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        ">
            <h2 style="color: white; margin: 0; text-align: center;">📚 Course Modules</h2>
        </div>
        """, unsafe_allow_html=True)
        
        modules = st.session_state.course_data['modules']
        
        # Module statistics dashboard
        total_items = sum(len(module.get('items', [])) for module in modules)
        file_items = sum(len([item for item in module.get('items', []) if item.get('type') == 'File']) for module in modules)
        assignment_items = sum(len([item for item in module.get('items', []) if item.get('type') == 'Assignment']) for module in modules)
        
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%);
                padding: 15px;
                border-radius: 8px;
                text-align: center;
                color: white;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            ">
                <h3 style="margin: 0; font-size: 24px;">📚</h3>
                <p style="margin: 5px 0 0 0; font-size: 18px; font-weight: bold;">{}</p>
                <p style="margin: 0; font-size: 12px; opacity: 0.9;">Total Modules</p>
            </div>
            """.format(len(modules)), unsafe_allow_html=True)
        
        with col2:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #fa709a 0%, #fee140 100%);
                padding: 15px;
                border-radius: 8px;
                text-align: center;
                color: white;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            ">
                <h3 style="margin: 0; font-size: 24px;">📋</h3>
                <p style="margin: 5px 0 0 0; font-size: 18px; font-weight: bold;">{}</p>
                <p style="margin: 0; font-size: 12px; opacity: 0.9;">Total Items</p>
            </div>
            """.format(total_items), unsafe_allow_html=True)
        
        with col3:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #a8edea 0%, #fed6e3 100%);
                padding: 15px;
                border-radius: 8px;
                text-align: center;
                color: #333;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            ">
                <h3 style="margin: 0; font-size: 24px;">📄</h3>
                <p style="margin: 5px 0 0 0; font-size: 18px; font-weight: bold;">{}</p>
                <p style="margin: 0; font-size: 12px; opacity: 0.8;">Files</p>
            </div>
            """.format(file_items), unsafe_allow_html=True)
        
        with col4:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #ffecd2 0%, #fcb69f 100%);
                padding: 15px;
                border-radius: 8px;
                text-align: center;
                color: #333;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            ">
                <h3 style="margin: 0; font-size: 24px;">📝</h3>
                <p style="margin: 5px 0 0 0; font-size: 18px; font-weight: bold;">{}</p>
                <p style="margin: 0; font-size: 12px; opacity: 0.8;">Assignments</p>
            </div>
            """.format(assignment_items), unsafe_allow_html=True)
        
        st.markdown("<br>", unsafe_allow_html=True)
        
        # Enhanced module display
        for i, module in enumerate(modules, 1):
            # Dynamic color coding based on module position
            colors = [
                "linear-gradient(135deg, #667eea 0%, #764ba2 100%)",
                "linear-gradient(135deg, #f093fb 0%, #f5576c 100%)",
                "linear-gradient(135deg, #4facfe 0%, #00f2fe 100%)",
                "linear-gradient(135deg, #43e97b 0%, #38f9d7 100%)",
                "linear-gradient(135deg, #fa709a 0%, #fee140 100%)"
            ]
            module_color = colors[(i-1) % len(colors)]
            
            # Module header with enhanced styling
            st.markdown(f"""
            <div style="
                background: {module_color};
                padding: 15px;
                border-radius: 10px 10px 0 0;
                margin-top: 20px;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            ">
                <h3 style="color: white; margin: 0;">📚 Module {i}: {module.get('name', 'Unnamed Module')}</h3>
                <p style="color: rgba(255,255,255,0.9); margin: 5px 0 0 0; font-size: 14px;">{len(module.get('items', []))} items</p>
            </div>
            """, unsafe_allow_html=True)
            
            with st.expander("📋 View Module Contents", expanded=False):
                # Module description and stats
                col1, col2 = st.columns([2, 1])
                with col1:
                    if module.get('description'):
                        st.markdown(f"📄 **Description:** {module.get('description', 'No description')}")
                    else:
                        st.markdown("📄 **Description:** *No description available*")
                
                with col2:
                    st.metric("📊 Items", len(module.get('items', [])))
                
                st.divider()
                
                # Show module items with enhanced formatting
                items = module.get('items', [])
                if items:
                    # Group items by type for better organization
                    item_types = {}
                    for item in items:
                        item_type = item.get('type', 'Unknown')
                        if item_type not in item_types:
                            item_types[item_type] = []
                        item_types[item_type].append(item)
                    
                    # Display items grouped by type with color coding
                    for item_type, type_items in item_types.items():
                        # Get appropriate icon and color for item type
                        type_config = {
                            'File': {'icon': '📄', 'color': '#3498db'},
                            'Page': {'icon': '📝', 'color': '#2ecc71'},
                            'Assignment': {'icon': '📝', 'color': '#e74c3c'},
                            'Discussion': {'icon': '💬', 'color': '#9b59b6'},
                            'Quiz': {'icon': '❓', 'color': '#f39c12'},
                            'ExternalUrl': {'icon': '🔗', 'color': '#1abc9c'},
                            'ExternalTool': {'icon': '🛠️', 'color': '#34495e'}
                        }
                        config = type_config.get(item_type, {'icon': '📌', 'color': '#95a5a6'})
                        
                        st.markdown(f"""
                        <div style="
                            background: {config['color']};
                            color: white;
                            padding: 10px;
                            border-radius: 5px;
                            margin: 10px 0;
                        ">
                            <h4 style="margin: 0;">{config['icon']} {item_type}s ({len(type_items)})</h4>
                        </div>
                        """, unsafe_allow_html=True)
                        
                        for j, item in enumerate(type_items, 1):
                            item_title = item.get('title', 'Untitled')
                            item_url = item.get('html_url', '')
                            
                            col1, col2 = st.columns([3, 1])
                            with col1:
                                if item_url:
                                    st.markdown(f"   {j}. [{item_title}]({item_url}) 🔗")
                                else:
                                    st.markdown(f"   {j}. {item_title}")
                            
                            with col2:
                                # Add download button for file items
                                if item_type == 'File' and item.get('content_id'):
                                    try:
                                        file_id = item.get('content_id')
                                        if st.button("📥 Download", key=f"download_module_{file_id}"):
                                             with st.spinner("Downloading file..."):
                                                 try:
                                                     if not st.session_state.get('client'):
                                                         st.error("Canvas client not available. Please log in again.")
                                                         return
                                                     file_content, filename, content_type = st.session_state.client.download_file_content(file_id)
                                                     st.download_button(
                                                         label="📥 Click to save file",
                                                         data=file_content,
                                                         file_name=filename,
                                                         mime=content_type,
                                                         key=f"save_module_{file_id}"
                                                     )
                                                 except Exception as e:
                                                     st.write("❌ Download unavailable")
                                    except Exception as e:
                                        st.write("❌ Download unavailable")
                else:
                    st.info("No items found in this module.")
    
    def render_files_section(self):
        """Render enhanced course files section with improved visual design and quick actions"""
        if not st.session_state.course_data.get('files'):
            return
        
        # Enhanced header with gradient background
        st.markdown("""
        <div style="
            background: linear-gradient(135deg, #ff9a9e 0%, #fecfef 50%, #fecfef 100%);
            padding: 20px;
            border-radius: 10px;
            margin-bottom: 20px;
            box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        ">
            <h2 style="color: white; margin: 0; text-align: center;">📁 Course Files</h2>
        </div>
        """, unsafe_allow_html=True)
        
        files = st.session_state.course_data['files']
        
        # Check if files are local
        local_files = [f for f in files if f.get('is_local', False)]
        api_files = [f for f in files if not f.get('is_local', False)]
        
        # File statistics dashboard
        total_size = sum(f.get('size', 0) for f in files)
        file_types_count = len(set([f.get('content-type', 'Unknown').split('/')[0] for f in files]))
        
        # Format total size
        if total_size > 1024 * 1024 * 1024:
            total_size_str = f"{total_size / (1024 * 1024 * 1024):.1f} GB"
        elif total_size > 1024 * 1024:
            total_size_str = f"{total_size / (1024 * 1024):.1f} MB"
        elif total_size > 1024:
            total_size_str = f"{total_size / 1024:.1f} KB"
        else:
            total_size_str = f"{total_size} B"
        
        # Show file source status
        # if local_files:
        #     st.success(f"📁 **{len(local_files)} files loaded from local downloads directory** (No Canvas API permissions needed)")
        # elif api_files:
        #     st.info(f"🌐 **{len(api_files)} files loaded from Canvas API**")
        # else:
        #     st.warning("⚠️ No files found")
        
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                padding: 15px;
                border-radius: 8px;
                text-align: center;
                color: white;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            ">
                <h3 style="margin: 0; font-size: 24px;">📁</h3>
                <p style="margin: 5px 0 0 0; font-size: 18px; font-weight: bold;">{}</p>
                <p style="margin: 0; font-size: 12px; opacity: 0.9;">Total Files</p>
            </div>
            """.format(len(files)), unsafe_allow_html=True)
        
        with col2:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);
                padding: 15px;
                border-radius: 8px;
                text-align: center;
                color: white;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            ">
                <h3 style="margin: 0; font-size: 24px;">💾</h3>
                <p style="margin: 5px 0 0 0; font-size: 18px; font-weight: bold;">{}</p>
                <p style="margin: 0; font-size: 12px; opacity: 0.9;">Total Size</p>
            </div>
            """.format(total_size_str), unsafe_allow_html=True)
        
        with col3:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%);
                padding: 15px;
                border-radius: 8px;
                text-align: center;
                color: white;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            ">
                <h3 style="margin: 0; font-size: 24px;">🗂️</h3>
                <p style="margin: 5px 0 0 0; font-size: 18px; font-weight: bold;">{}</p>
                <p style="margin: 0; font-size: 12px; opacity: 0.9;">File Types</p>
            </div>
            """.format(file_types_count), unsafe_allow_html=True)
        
        with col4:
            # Quick actions for files
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #43e97b 0%, #38f9d7 100%);
                padding: 15px;
                border-radius: 8px;
                text-align: center;
                color: white;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            ">
                <h3 style="margin: 0; font-size: 24px;">🚀</h3>
                <p style="margin: 5px 0 0 0; font-size: 14px; font-weight: bold;">Quick Actions</p>
            </div>
            """, unsafe_allow_html=True)
        
        st.markdown("<br>", unsafe_allow_html=True)
        
        # Quick action buttons
        st.markdown("### 🚀 Quick Actions")
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            if st.button("📥 Download All Files", type="primary", use_container_width=True):
                st.info("Bulk download feature - would download all files")
        
        with col2:
            if st.button("📊 Export File List", type="secondary", use_container_width=True):
                st.info("Export file metadata to CSV")
        
        with col3:
            if st.button("🔍 Advanced Search", type="secondary", use_container_width=True):
                st.session_state.show_advanced_search = not st.session_state.get('show_advanced_search', False)
                st.rerun()
        
        with col4:
            if st.button("📁 Open Downloads", type="secondary", use_container_width=True):
                st.session_state.active_tab = "download"
                st.rerun()
        
        st.markdown("<br>", unsafe_allow_html=True)
        
        # Enhanced file filters
        st.markdown("### 🔍 Filter Files")
        col1, col2, col3 = st.columns([1, 1, 1])
        
        with col1:
            # Filter by file type
            file_types = list(set([f.get('content-type', 'Unknown').split('/')[0] for f in files]))
            selected_type = st.selectbox("📂 Filter by type:", ['All'] + file_types)
        
        with col2:
            # Search files
            search_term = st.text_input("🔍 Search files:", placeholder="Enter filename...")
        
        with col3:
            # Size filter
            size_filter = st.selectbox("📏 Size filter:", ['All sizes', 'Small (<1MB)', 'Medium (1-10MB)', 'Large (>10MB)'])
        
        # Advanced search (if enabled)
        if st.session_state.get('show_advanced_search', False):
            st.markdown("#### 🔍 Advanced Search Options")
            col1, col2 = st.columns(2)
            with col1:
                date_filter = st.selectbox("📅 Date filter:", ['All dates', 'Last week', 'Last month', 'Last semester'])
            with col2:
                sort_by = st.selectbox("📊 Sort by:", ['Name', 'Size', 'Type', 'Date modified'])
        
        # Filter files
        filtered_files = files.copy()
        
        if selected_type != 'All':
            filtered_files = [f for f in filtered_files if f.get('content-type', '').startswith(selected_type)]
        
        if search_term:
            filtered_files = [f for f in filtered_files if search_term.lower() in f.get('display_name', '').lower()]
        
        # Size filtering
        if size_filter != 'All sizes':
            if size_filter == 'Small (<1MB)':
                filtered_files = [f for f in filtered_files if f.get('size', 0) < 1024 * 1024]
            elif size_filter == 'Medium (1-10MB)':
                filtered_files = [f for f in filtered_files if 1024 * 1024 <= f.get('size', 0) <= 10 * 1024 * 1024]
            elif size_filter == 'Large (>10MB)':
                filtered_files = [f for f in filtered_files if f.get('size', 0) > 10 * 1024 * 1024]
        
        # Display files with enhanced cards
        if filtered_files:
            st.markdown(f"### 📋 Found {len(filtered_files)} files")
            
            for file_info in filtered_files:
                # Enhanced file card with gradient background
                file_name = file_info.get('display_name', 'Unknown')
                file_size = file_info.get('size', 0)
                file_type = file_info.get('content-type', 'Unknown')
                
                # Format file size
                if file_size > 1024 * 1024:
                    size_str = f"{file_size / (1024 * 1024):.1f} MB"
                elif file_size > 1024:
                    size_str = f"{file_size / 1024:.1f} KB"
                else:
                    size_str = f"{file_size} B"
                
                # File type icon
                type_icons = {
                    'application': '📄',
                    'image': '🖼️',
                    'video': '🎥',
                    'audio': '🎵',
                    'text': '📝'
                }
                file_icon = type_icons.get(file_type.split('/')[0], '📎')
                
                st.markdown(f"""
                <div style="
                    background: linear-gradient(135deg, #f8f9fa 0%, #e9ecef 100%);
                    padding: 15px;
                    border-radius: 10px;
                    margin: 10px 0;
                    border-left: 4px solid #007bff;
                    box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
                ">
                    <div style="display: flex; align-items: center; margin-bottom: 8px;">
                        <span style="font-size: 24px; margin-right: 10px;">{file_icon}</span>
                        <h4 style="margin: 0; color: #333;">{file_name}</h4>
                    </div>
                    <p style="margin: 0; color: #666; font-size: 14px;">Type: {file_type} | Size: {size_str}</p>
                </div>
                """, unsafe_allow_html=True)
                
                # Action buttons
                col1, col2, col3, col4 = st.columns([2, 1, 1, 1])
                
                with col2:
                    # View button
                    if file_info.get('url'):
                        st.markdown(f"[👁️ View]({file_info.get('url')})")
                    else:
                        st.write("👁️ View")
                
                with col3:
                    # Download button
                    try:
                        file_id = file_info.get('id')
                        if file_id:
                            if st.button("📥 Download", key=f"download_file_{file_id}", use_container_width=True):
                                with st.spinner("Downloading file..."):
                                    if not st.session_state.get('client'):
                                        st.error("Canvas client not available. Please log in again.")
                                        return
                                    file_content, filename, content_type = st.session_state.client.download_file_content(file_id)
                                    st.download_button(
                                        label="📥 Click to save file",
                                        data=file_content,
                                        file_name=filename,
                                        mime=content_type,
                                        key=f"save_file_{file_id}"
                                    )
                        else:
                            st.button("❌ Unavailable", disabled=True, use_container_width=True)
                    except Exception as e:
                        st.button("❌ Error", disabled=True, use_container_width=True)
                
                with col4:
                    # Share/Info button
                    if st.button("ℹ️ Info", key=f"info_file_{file_info.get('id', 'unknown')}", use_container_width=True):
                        st.info(f"File ID: {file_info.get('id')}\nURL: {file_info.get('url', 'N/A')}")
        else:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #ffeaa7 0%, #fab1a0 100%);
                padding: 20px;
                border-radius: 10px;
                text-align: center;
                color: #2d3436;
            ">
                <h3 style="margin: 0;">🔍 No files found</h3>
                <p style="margin: 10px 0 0 0;">Try adjusting your search criteria or filters</p>
            </div>
            """, unsafe_allow_html=True)
    
    def render_export_section(self):
        """Render enhanced data export section with comprehensive options and statistics"""
        # Enhanced header with gradient background
        st.markdown("""
        <div style="
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            padding: 20px;
            border-radius: 10px;
            margin-bottom: 20px;
            box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        ">
            <h2 style="color: white; margin: 0; text-align: center;">📊 Export Data</h2>
            <p style="color: rgba(255,255,255,0.9); margin: 10px 0 0 0; text-align: center;">Export your Canvas data in various formats</p>
        </div>
        """, unsafe_allow_html=True)
        
        # Export statistics dashboard
        if hasattr(self, 'export_dir') and self.export_dir.exists():
            export_files = list(self.export_dir.glob("*.json"))
            csv_files = list(self.export_dir.glob("*.csv"))
            total_exports = len(export_files) + len(csv_files)
            
            # Calculate total export size
            total_size = sum(f.stat().st_size for f in export_files + csv_files if f.exists())
            if total_size > 1024 * 1024:
                size_str = f"{total_size / (1024 * 1024):.1f} MB"
            elif total_size > 1024:
                size_str = f"{total_size / 1024:.1f} KB"
            else:
                size_str = f"{total_size} B"
        else:
            export_files = []
            csv_files = []
            total_exports = 0
            size_str = "0 B"
        
        # Statistics cards
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                padding: 15px;
                border-radius: 8px;
                text-align: center;
                color: white;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            ">
                <h3 style="margin: 0; font-size: 24px;">📦</h3>
                <p style="margin: 5px 0 0 0; font-size: 18px; font-weight: bold;">{}</p>
                <p style="margin: 0; font-size: 12px; opacity: 0.9;">Total Exports</p>
            </div>
            """.format(total_exports), unsafe_allow_html=True)
        
        with col2:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);
                padding: 15px;
                border-radius: 8px;
                text-align: center;
                color: white;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            ">
                <h3 style="margin: 0; font-size: 24px;">📄</h3>
                <p style="margin: 5px 0 0 0; font-size: 18px; font-weight: bold;">{}</p>
                <p style="margin: 0; font-size: 12px; opacity: 0.9;">JSON Files</p>
            </div>
            """.format(len(export_files)), unsafe_allow_html=True)
        
        with col3:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%);
                padding: 15px;
                border-radius: 8px;
                text-align: center;
                color: white;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            ">
                <h3 style="margin: 0; font-size: 24px;">📊</h3>
                <p style="margin: 5px 0 0 0; font-size: 18px; font-weight: bold;">{}</p>
                <p style="margin: 0; font-size: 12px; opacity: 0.9;">CSV Files</p>
            </div>
            """.format(len(csv_files)), unsafe_allow_html=True)
        
        with col4:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #43e97b 0%, #38f9d7 100%);
                padding: 15px;
                border-radius: 8px;
                text-align: center;
                color: white;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            ">
                <h3 style="margin: 0; font-size: 24px;">💾</h3>
                <p style="margin: 5px 0 0 0; font-size: 18px; font-weight: bold;">{}</p>
                <p style="margin: 0; font-size: 12px; opacity: 0.9;">Total Size</p>
            </div>
            """.format(size_str), unsafe_allow_html=True)
        
        st.markdown("<br>", unsafe_allow_html=True)
        
        # Quick export actions
        st.markdown("### 🚀 Quick Export Actions")
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            if st.button("📦 Export All Courses", type="primary", use_container_width=True):
                self.batch_export_all_courses()
        
        with col2:
            if st.button("📁 Open Export Folder", type="secondary", use_container_width=True):
                st.info(f"Export folder: {self.export_dir if hasattr(self, 'export_dir') else 'Not set'}")
        
        with col3:
            if st.button("🗑️ Clean Old Exports", type="secondary", use_container_width=True):
                st.info("Clean exports older than 30 days")
        
        with col4:
            if st.button("📊 Export Analytics", type="secondary", use_container_width=True):
                st.info("Generate export usage analytics")
        
        st.markdown("<br>", unsafe_allow_html=True)
        
        # Enhanced batch export section
        st.markdown("### 🚀 Batch Export Configuration")
        st.markdown("""
        <div style="
            background: linear-gradient(135deg, #f8f9fa 0%, #e9ecef 100%);
            padding: 15px;
            border-radius: 10px;
            border-left: 4px solid #28a745;
            margin: 10px 0;
        ">
            <p style="margin: 0; color: #333;">📋 Export comprehensive JSON data for all your courses at once with advanced options.</p>
        </div>
        """, unsafe_allow_html=True)
        
        col1, col2, col3 = st.columns([1, 1, 1])
        
        with col1:
            export_format = st.selectbox(
                "📦 Export Format",
                ["Individual JSON files", "Single ZIP file", "Structured folders"],
                help="Choose how to package the exported data"
            )
        
        with col2:
            include_options = st.multiselect(
                "📋 Include Data",
                ["File metadata", "Assignment details", "Module content", "Grades", "Discussions"],
                default=["File metadata", "Assignment details", "Module content"],
                help="Select what data to include in exports"
            )
        
        with col3:
            export_quality = st.selectbox(
                "🎯 Export Quality",
                ["Basic", "Standard", "Comprehensive"],
                index=1,
                help="Choose the level of detail in exports"
            )
        
        # Advanced export options
        with st.expander("⚙️ Advanced Export Options"):
            col1, col2 = st.columns(2)
            
            with col1:
                date_range = st.selectbox(
                    "📅 Date Range",
                    ["All time", "Current semester", "Last 6 months", "Last year"],
                    help="Filter exports by date range"
                )
                
                compression = st.selectbox(
                    "🗜️ Compression",
                    ["None", "ZIP", "7Z"],
                    index=1,
                    help="Choose compression method"
                )
            
            with col2:
                max_file_size = st.slider(
                    "📏 Max File Size (MB)",
                    min_value=1,
                    max_value=100,
                    value=50,
                    help="Maximum size for individual export files"
                )
                
                parallel_exports = st.slider(
                    "⚡ Parallel Exports",
                    min_value=1,
                    max_value=5,
                    value=2,
                    help="Number of courses to export simultaneously"
                )
        
        st.markdown("<br>", unsafe_allow_html=True)
        
        # Current course export section
        if st.session_state.course_data:
            st.markdown("### 📄 Current Course Export")
            
            course = st.session_state.selected_course
            course_name = course.get('name', 'Unknown') if course else 'Unknown'
            course_data = st.session_state.course_data
            
            # Course export statistics
            assignments_count = len(course_data.get('assignments', []))
            modules_count = len(course_data.get('modules', []))
            files_count = len(course_data.get('files', []))
            
            st.markdown(f"""
            <div style="
                background: linear-gradient(135deg, #e3f2fd 0%, #bbdefb 100%);
                padding: 15px;
                border-radius: 10px;
                margin: 10px 0;
                border-left: 4px solid #2196f3;
            ">
                <h4 style="margin: 0 0 10px 0; color: #1976d2;">📚 {course_name}</h4>
                <div style="display: flex; gap: 20px; flex-wrap: wrap;">
                    <span style="color: #666;">📝 {assignments_count} assignments</span>
                    <span style="color: #666;">📂 {modules_count} modules</span>
                    <span style="color: #666;">📎 {files_count} files</span>
                </div>
            </div>
            """, unsafe_allow_html=True)
            
            # Export options with enhanced buttons
            col1, col2, col3, col4 = st.columns(4)
            
            with col1:
                if st.button("📄 Export to JSON", type="primary", use_container_width=True):
                    self.export_to_json(course_data, course_name, course.get('id') if course else None)
            
            with col2:
                if st.button("📊 Export Assignments CSV", type="secondary", use_container_width=True):
                    if course_data.get('assignments'):
                        self.export_assignments_to_csv(course_data['assignments'], course_name)
                    else:
                        st.warning("No assignments to export")
            
            with col3:
                if st.button("📁 Export Files List", type="secondary", use_container_width=True):
                    if course_data.get('files'):
                        st.info("Export file metadata to CSV")
                    else:
                        st.warning("No files to export")
            
            with col4:
                if st.button("📋 Export Summary", type="secondary", use_container_width=True):
                    st.info("Generate course summary report")
        else:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #fff3cd 0%, #ffeaa7 100%);
                padding: 20px;
                border-radius: 10px;
                text-align: center;
                color: #856404;
                border-left: 4px solid #ffc107;
            ">
                <h4 style="margin: 0;">📚 No Course Selected</h4>
                <p style="margin: 10px 0 0 0;">Select a course from the Overview tab to enable single course export options</p>
            </div>
            """, unsafe_allow_html=True)
    
    def render_download_section(self):
        """Render enhanced download section for downloading files from JSON exports with comprehensive monitoring"""
        # Enhanced header with gradient background
        st.markdown("""
        <div style="
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            padding: 20px;
            border-radius: 10px;
            margin-bottom: 20px;
            box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        ">
            <h2 style="color: white; margin: 0; text-align: center;">⬇️ Download Files from JSON Exports</h2>
            <p style="color: rgba(255,255,255,0.9); margin: 10px 0 0 0; text-align: center;">Upload JSON export files and download course materials with advanced monitoring</p>
        </div>
        """, unsafe_allow_html=True)
        
        # Initialize session state for downloads
        if 'download_progress' not in st.session_state:
            st.session_state.download_progress = []
        if 'download_stats' not in st.session_state:
            st.session_state.download_stats = {}
        if 'download_history' not in st.session_state:
            st.session_state.download_history = []
        
        # Download statistics dashboard
        total_downloads = len(st.session_state.download_history)
        current_stats = st.session_state.download_stats
        active_downloads = 1 if st.session_state.download_progress else 0
        
        # Statistics cards
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                padding: 15px;
                border-radius: 8px;
                text-align: center;
                color: white;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            ">
                <h3 style="margin: 0; font-size: 24px;">📥</h3>
                <p style="margin: 5px 0 0 0; font-size: 18px; font-weight: bold;">{}</p>
                <p style="margin: 0; font-size: 12px; opacity: 0.9;">Total Sessions</p>
            </div>
            """.format(total_downloads), unsafe_allow_html=True)
        
        with col2:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #43e97b 0%, #38f9d7 100%);
                padding: 15px;
                border-radius: 8px;
                text-align: center;
                color: white;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            ">
                <h3 style="margin: 0; font-size: 24px;">✅</h3>
                <p style="margin: 5px 0 0 0; font-size: 18px; font-weight: bold;">{}</p>
                <p style="margin: 0; font-size: 12px; opacity: 0.9;">Downloaded</p>
            </div>
            """.format(current_stats.get('downloaded', 0)), unsafe_allow_html=True)
        
        with col3:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%);
                padding: 15px;
                border-radius: 8px;
                text-align: center;
                color: white;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            ">
                <h3 style="margin: 0; font-size: 24px;">⏭️</h3>
                <p style="margin: 5px 0 0 0; font-size: 18px; font-weight: bold;">{}</p>
                <p style="margin: 0; font-size: 12px; opacity: 0.9;">Skipped</p>
            </div>
            """.format(current_stats.get('skipped', 0)), unsafe_allow_html=True)
        
        with col4:
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);
                padding: 15px;
                border-radius: 8px;
                text-align: center;
                color: white;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            ">
                <h3 style="margin: 0; font-size: 24px;">❌</h3>
                <p style="margin: 5px 0 0 0; font-size: 18px; font-weight: bold;">{}</p>
                <p style="margin: 0; font-size: 12px; opacity: 0.9;">Errors</p>
            </div>
            """.format(current_stats.get('errors', 0)), unsafe_allow_html=True)
        
        st.markdown("<br>", unsafe_allow_html=True)
        
        # Quick action buttons
        st.markdown("### 🚀 Quick Actions")
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            if st.button("📁 Open Downloads", type="secondary", use_container_width=True):
                st.info("Open downloads folder")
        
        with col2:
            if st.button("🔄 Refresh Exports", type="secondary", use_container_width=True):
                st.rerun()
        
        with col3:
            if st.button("📊 View History", type="secondary", use_container_width=True):
                st.info("Show download history")
        
        with col4:
            if st.button("🗑️ Clear All", type="secondary", use_container_width=True):
                st.session_state.download_progress = []
                st.session_state.download_stats = {}
                st.success("Cleared all progress")
        
        st.markdown("<br>", unsafe_allow_html=True)
        
        # Enhanced file upload section
        st.markdown("### 📁 Select JSON Files")
        st.markdown("""
        <div style="
            background: linear-gradient(135deg, #f8f9fa 0%, #e9ecef 100%);
            padding: 15px;
            border-radius: 10px;
            border-left: 4px solid #28a745;
            margin: 10px 0;
        ">
            <p style="margin: 0; color: #333;">📋 Upload JSON export files or select from existing exports to download course materials.</p>
        </div>
        """, unsafe_allow_html=True)
        
        # Option 1: Upload files
        uploaded_files = st.file_uploader(
            "📤 Upload JSON export files",
            type=['json'],
            accept_multiple_files=True,
            help="Upload one or more JSON export files from the Export tab"
        )
        
        # Option 2: Select from exports directory
        export_files = []
        if self.export_dir.exists():
            export_files = list(self.export_dir.glob("*.json"))
        
        if export_files:
            st.markdown("### 📂 Available Export Files")
            
            # Show export files with details
            for i, file_path in enumerate(export_files[:5]):  # Show first 5 files
                file_size = file_path.stat().st_size
                if file_size > 1024 * 1024:
                    size_str = f"{file_size / (1024 * 1024):.1f} MB"
                elif file_size > 1024:
                    size_str = f"{file_size / 1024:.1f} KB"
                else:
                    size_str = f"{file_size} B"
                
                st.markdown(f"""
                <div style="
                    background: white;
                    padding: 10px;
                    border-radius: 8px;
                    margin: 5px 0;
                    border-left: 4px solid #007bff;
                    box-shadow: 0 1px 3px rgba(0, 0, 0, 0.1);
                ">
                    <div style="display: flex; justify-content: space-between; align-items: center;">
                        <span style="font-weight: bold; color: #333;">{file_path.name}</span>
                        <span style="color: #666; font-size: 12px;">{size_str}</span>
                    </div>
                </div>
                """, unsafe_allow_html=True)
            
            if len(export_files) > 5:
                st.info(f"... and {len(export_files) - 5} more files")
            
            selected_export_files = st.multiselect(
                "📋 Select files to process:",
                options=export_files,
                format_func=lambda x: f"{x.name} ({x.stat().st_size / (1024*1024):.1f} MB)"
            )
        else:
            selected_export_files = []
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #fff3cd 0%, #ffeaa7 100%);
                padding: 15px;
                border-radius: 10px;
                text-align: center;
                color: #856404;
                border-left: 4px solid #ffc107;
            ">
                <h4 style="margin: 0;">📂 No Export Files Found</h4>
                <p style="margin: 10px 0 0 0;">Create exports first from the Export tab</p>
            </div>
            """, unsafe_allow_html=True)
        
        # Enhanced download options
        st.markdown("### ⚙️ Download Configuration")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("**🔧 Basic Options**")
            dry_run = st.checkbox(
                "🔍 Dry Run Mode",
                value=False,
                help="Preview what would be downloaded without actually downloading"
            )
            
            use_api_token = st.checkbox(
                "🔑 Use API Token for Downloads",
                value=True,
                help="Use your Canvas API token for authenticated downloads"
            )
            
            download_dir = st.text_input(
                "📁 Download Directory",
                value="downloads",
                help="Directory where files will be saved"
            )
        
        with col2:
            st.markdown("**⚡ Performance Options**")
            max_concurrent = st.slider(
                "🚀 Max Concurrent Downloads",
                min_value=1,
                max_value=10,
                value=3,
                help="Number of files to download simultaneously"
            )
            
            timeout = st.slider(
                "⏱️ Timeout (seconds)",
                min_value=10,
                max_value=300,
                value=60,
                help="Timeout for individual file downloads"
            )
            
            retry_attempts = st.slider(
                "🔄 Retry Attempts",
                min_value=1,
                max_value=5,
                value=3,
                help="Number of retry attempts for failed downloads"
            )
        
        # Advanced options
        with st.expander("🔧 Advanced Download Options"):
            col1, col2 = st.columns(2)
            
            with col1:
                file_types = st.multiselect(
                    "📎 File Types to Download",
                    ["PDF", "DOCX", "PPTX", "Images", "Videos", "Audio", "Archives", "All"],
                    default=["All"],
                    help="Select specific file types to download"
                )
                
                min_file_size = st.number_input(
                    "📏 Min File Size (KB)",
                    min_value=0,
                    value=0,
                    help="Skip files smaller than this size"
                )
            
            with col2:
                max_file_size = st.number_input(
                    "📐 Max File Size (MB)",
                    min_value=1,
                    value=100,
                    help="Skip files larger than this size"
                )
                
                organize_by = st.selectbox(
                    "📂 Organize Files By",
                    ["Course", "File Type", "Date", "Flat Structure"],
                    help="How to organize downloaded files"
                )
        
        # Process files
        files_to_process = []
        
        # Add uploaded files
        if uploaded_files:
            for uploaded_file in uploaded_files:
                # Save uploaded file temporarily
                temp_path = self.export_dir / uploaded_file.name
                with open(temp_path, 'wb') as f:
                    f.write(uploaded_file.getbuffer())
                files_to_process.append(temp_path)
        
        # Add selected export files
        files_to_process.extend(selected_export_files)
        
        if files_to_process:
            st.markdown("### 📋 Files Ready for Processing")
            
            # Show files in a nice format
            total_size = 0
            for file_path in files_to_process:
                file_size = file_path.stat().st_size
                total_size += file_size
                
                if file_size > 1024 * 1024:
                    size_str = f"{file_size / (1024 * 1024):.1f} MB"
                elif file_size > 1024:
                    size_str = f"{file_size / 1024:.1f} KB"
                else:
                    size_str = f"{file_size} B"
                
                st.markdown(f"""
                <div style="
                    background: linear-gradient(135deg, #e3f2fd 0%, #bbdefb 100%);
                    padding: 10px;
                    border-radius: 8px;
                    margin: 5px 0;
                    border-left: 4px solid #2196f3;
                ">
                    <div style="display: flex; justify-content: space-between; align-items: center;">
                        <span style="font-weight: bold; color: #1976d2;">📄 {file_path.name}</span>
                        <span style="color: #666; font-size: 12px;">{size_str}</span>
                    </div>
                </div>
                """, unsafe_allow_html=True)
            
            # Summary
            total_size_str = f"{total_size / (1024 * 1024):.1f} MB" if total_size > 1024 * 1024 else f"{total_size / 1024:.1f} KB"
            st.info(f"📊 Ready to process {len(files_to_process)} files ({total_size_str} total)")
            
            # Download button
            col1, col2 = st.columns([3, 1])
            with col1:
                if st.button("🚀 Start Download Process", type="primary", use_container_width=True):
                    self.start_download_process(files_to_process, download_dir, dry_run, use_api_token)
            with col2:
                if st.button("🔍 Preview Only", type="secondary", use_container_width=True):
                    self.start_download_process(files_to_process, download_dir, True, use_api_token)
        
        # Enhanced progress display
        if st.session_state.download_progress:
            st.markdown("### 📊 Download Progress & Monitoring")
            
            # Progress container with enhanced styling
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #f8f9fa 0%, #e9ecef 100%);
                padding: 15px;
                border-radius: 10px;
                margin: 10px 0;
                border-left: 4px solid #28a745;
            ">
            """, unsafe_allow_html=True)
            
            # Show recent progress messages in a scrollable container
            with st.container():
                st.markdown("**📝 Recent Activity:**")
                
                # Create a text area for progress messages
                progress_text = "\n".join(st.session_state.download_progress[-20:])
                st.text_area(
                    "Progress Log",
                    value=progress_text,
                    height=200,
                    disabled=True,
                    label_visibility="collapsed"
                )
                
                # Enhanced statistics display
                if st.session_state.download_stats:
                    stats = st.session_state.download_stats
                    
                    st.markdown("**📈 Current Session Statistics:**")
                    col1, col2, col3, col4, col5 = st.columns(5)
                    
                    with col1:
                        st.metric("📁 Total Files", stats.get('total_files', 0))
                    
                    with col2:
                        st.metric("✅ Downloaded", stats.get('downloaded', 0))
                    
                    with col3:
                        st.metric("⏭️ Skipped", stats.get('skipped', 0))
                    
                    with col4:
                        st.metric("❌ Errors", stats.get('errors', 0))
                    
                    with col5:
                        success_rate = 0
                        if stats.get('total_files', 0) > 0:
                            success_rate = (stats.get('downloaded', 0) / stats.get('total_files', 1)) * 100
                        st.metric("📊 Success Rate", f"{success_rate:.1f}%")
            
            st.markdown("</div>", unsafe_allow_html=True)
        
        # Action buttons for progress management
        if st.session_state.download_progress:
            col1, col2, col3 = st.columns(3)
            
            with col1:
                if st.button("🗑️ Clear Progress Log", use_container_width=True):
                    st.session_state.download_progress = []
                    st.rerun()
            
            with col2:
                if st.button("📋 Export Log", use_container_width=True):
                    st.info("Export progress log to file")
            
            with col3:
                if st.button("🔄 Refresh Status", use_container_width=True):
                    st.rerun()
    
    def start_download_process(self, files_to_process: List[Path], download_dir: str, dry_run: bool, use_api_token: bool):
        """Start the download process for selected files"""
        try:
            # Clear previous progress
            st.session_state.download_progress = []
            st.session_state.download_stats = {}
            
            # Get API token if requested
            api_token = None
            base_url = None
            
            if use_api_token and self.client:
                api_token = self.client.api_token
                base_url = self.client.base_url
            
            # Initialize downloader
            downloader = CanvasFileDownloader(
                api_token=api_token,
                base_url=base_url,
                output_dir=download_dir,
                dry_run=dry_run
            )
            
            # Progress callback function
            def progress_callback(message: str):
                st.session_state.download_progress.append(f"[{datetime.now().strftime('%H:%M:%S')}] {message}")
                # Keep only last 100 messages to prevent memory issues
                if len(st.session_state.download_progress) > 100:
                    st.session_state.download_progress = st.session_state.download_progress[-100:]
            
            progress_callback(f"Starting download process for {len(files_to_process)} files...")
            progress_callback(f"Mode: {'Dry Run' if dry_run else 'Download'}")
            progress_callback(f"Output Directory: {download_dir}")
            
            # Process each JSON file
            for json_file in files_to_process:
                progress_callback(f"Processing: {json_file.name}")
                downloader.process_json_export(str(json_file), progress_callback)
            
            # Update statistics
            st.session_state.download_stats = downloader.stats.copy()
            
            progress_callback("Download process completed!")
            progress_callback(f"Summary: {downloader.stats['downloaded']} downloaded, {downloader.stats['skipped']} skipped, {downloader.stats['errors']} errors")
            
            # Show completion message
            if dry_run:
                st.success(f"✅ Dry run completed! Would have processed {downloader.stats['total_files']} files.")
            else:
                st.success(f"✅ Download completed! {downloader.stats['downloaded']} files downloaded.")
            
            # Rerun to update the display
            st.rerun()
            
        except Exception as e:
            error_msg = f"Error during download process: {str(e)}"
            st.session_state.download_progress.append(f"[{datetime.now().strftime('%H:%M:%S')}] ERROR: {error_msg}")
            st.error(error_msg)
            logging.error(f"Download process error: {e}", exc_info=True)
    
    def batch_export_all_courses(self):
        """Export all courses data to JSON files"""
        if not st.session_state.authenticated or not st.session_state.client:
            st.error("Please authenticate first")
            return
        
        try:
            with st.spinner("Exporting all courses data..."):
                # Get all courses if not already loaded
                if not st.session_state.courses:
                    st.session_state.courses = st.session_state.client.get_courses()
                
                export_files = []
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                total_courses = len(st.session_state.courses)
                
                for i, course in enumerate(st.session_state.courses):
                    course_id = course.get('id')
                    course_name = course.get('name', f'Course_{course_id}')
                    
                    status_text.text(f"Exporting: {course_name}")
                    
                    try:
                        # Use the enhanced export method from canvas_client
                        course_data = st.session_state.client.export_course_data(course_id)
                        
                        # Save to file
                        safe_name = "".join(c for c in course_name if c.isalnum() or c in (' ', '-', '_')).rstrip()
                        safe_name = safe_name.replace(' ', '_')
                        filename = f"course_{course_id}_{safe_name}.json"
                        filepath = self.export_dir / filename
                        
                        with open(filepath, 'w', encoding='utf-8') as f:
                            json.dump(course_data, f, indent=2, ensure_ascii=False)
                        
                        export_files.append(filepath)
                        logger.info(f"Exported course: {course_name}")
                        
                    except Exception as e:
                        logger.error(f"Failed to export course {course_name}: {e}")
                        st.warning(f"⚠️ Failed to export: {course_name}")
                    
                    # Update progress
                    progress_bar.progress((i + 1) / total_courses)
                
                status_text.text("Export completed!")
                
                if export_files:
                    st.success(f"✅ Successfully exported {len(export_files)} courses")
                    
                    # Create ZIP file for download
                    zip_buffer = io.BytesIO()
                    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                        for filepath in export_files:
                            zip_file.write(filepath, filepath.name)
                    
                    zip_buffer.seek(0)
                    
                    # Provide download button
                    st.download_button(
                        label="📥 Download All Exports (ZIP)",
                        data=zip_buffer.getvalue(),
                        file_name=f"canvas_courses_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip",
                        mime="application/zip"
                    )
                    
                    # Show individual files
                    with st.expander("📁 Individual Export Files"):
                        for filepath in export_files:
                            col1, col2 = st.columns([3, 1])
                            with col1:
                                st.write(f"📄 {filepath.name}")
                            with col2:
                                with open(filepath, 'rb') as f:
                                    st.download_button(
                                        label="📥",
                                        data=f.read(),
                                        file_name=filepath.name,
                                        mime="application/json",
                                        key=f"download_{filepath.stem}"
                                    )
                else:
                    st.error("❌ No courses were successfully exported")
                    
        except Exception as e:
            st.error(f"❌ Batch export failed: {str(e)}")
            logger.error(f"Batch export error: {e}")
    
    def export_to_json(self, data: Dict[str, Any], course_name: str, course_id: int = None):
        """Export data to JSON file"""
        try:
            safe_name = "".join(c for c in course_name if c.isalnum() or c in (' ', '-', '_')).rstrip()
            safe_name = safe_name.replace(' ', '_')
            
            if course_id:
                filename = f"course_{course_id}_{safe_name}.json"
            else:
                filename = f"{safe_name}_course_data.json"
            
            filepath = self.export_dir / filename
            
            # Add metadata to the export
            export_data = {
                'export_info': {
                    'exported_at': datetime.now().isoformat(),
                    'course_name': course_name,
                    'course_id': course_id,
                    'export_version': '2.0'
                },
                'course_data': data
            }
            
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(export_data, f, indent=2, ensure_ascii=False)
            
            st.success(f"✅ Data exported to {filepath}")
            
            # Provide download link
            with open(filepath, 'rb') as f:
                st.download_button(
                    label="📥 Download JSON File",
                    data=f.read(),
                    file_name=filename,
                    mime="application/json"
                )
                
        except Exception as e:
            st.error(f"❌ Export failed: {str(e)}")
    
    def export_assignments_to_csv(self, assignments: List[Dict], course_name: str):
        """Export assignments to CSV"""
        try:
            safe_name = "".join(c for c in course_name if c.isalnum() or c in (' ', '-', '_')).rstrip()
            safe_name = safe_name.replace(' ', '_')
            
            filename = f"{safe_name}_assignments.csv"
            filepath = self.export_dir / filename
            
            # Convert to DataFrame
            df = pd.DataFrame(assignments)
            
            # Select relevant columns
            columns_to_keep = ['name', 'due_at', 'points_possible', 'description', 'created_at']
            available_columns = [col for col in columns_to_keep if col in df.columns]
            
            if available_columns:
                df_export = df[available_columns]
                df_export.to_csv(filepath, index=False)
                
                st.success(f"✅ Assignments exported to {filepath}")
                
                # Provide download link
                with open(filepath, 'rb') as f:
                    st.download_button(
                        label="📥 Download CSV File",
                        data=f.read(),
                        file_name=filename,
                        mime="text/csv"
                    )
            else:
                st.warning("No suitable columns found for export")
                
        except Exception as e:
            st.error(f"❌ Export failed: {str(e)}")
    
    def render_ai_notes_section(self):
        """Render AI Notes section using exact LMS AI Assistant architecture"""
        st.header("🤖 AI-Powered Note Generation")
        st.markdown("Generate comprehensive lecture-style notes from your course documents using the exact same architecture as LMS AI Assistant.")
        
        # Import the exact notes generation functions
        try:
            from utils.notes import generate_notes_from_text, iter_generate_notes_from_texts
            from utils.text_processing import chunk_text
        except ImportError as e:
            st.error(f"❌ Failed to import notes utilities: {e}")
            st.info("💡 Make sure the utils directory exists with notes.py and text_processing.py")
            return
        
        # Initialize session state
        if 'ai_notes_cache' not in st.session_state:
            st.session_state.ai_notes_cache = {}
        if 'selected_file_for_ai' not in st.session_state:
            st.session_state.selected_file_for_ai = None
        
        # Two-column layout
        col1, col2 = st.columns([1, 2])
        
        with col1:
            st.subheader("📁 Course Files")
            self.render_course_file_selector_lms_style()
        
        with col2:
            st.subheader("🤖 AI Processing")
            self.render_ai_processing_panel_lms_style()
    
    def render_course_file_selector_lms_style(self):
        """Render course file selector using LMS AI Assistant style"""
        downloads_dir = os.path.join(os.path.dirname(__file__), 'downloads')
        
        if not os.path.exists(downloads_dir):
            st.warning("📁 No downloads folder found. Please download some course files first.")
            return
        
        # Get all course folders
        course_folders = [f for f in os.listdir(downloads_dir) 
                         if os.path.isdir(os.path.join(downloads_dir, f))]
        
        if not course_folders:
            st.warning("📁 No course folders found in downloads.")
            return
        
        # Course selection with expandable view
        for course_idx, course_folder in enumerate(course_folders):
            course_path = os.path.join(downloads_dir, course_folder)
            
            with st.expander(f"📚 {course_folder}", expanded=False):
                # Get files in this course folder
                try:
                    files = [f for f in os.listdir(course_path) 
                            if os.path.isfile(os.path.join(course_path, f)) and 
                            f.lower().endswith(('.pdf', '.docx', '.txt'))]
                    
                    if files:
                        st.markdown(f"**Found {len(files)} supported files:**")
                        
                        for file_idx, file_name in enumerate(sorted(files)):
                            file_path = os.path.join(course_path, file_name)
                            
                            # File selection button with unique key
                            unique_key = f"lms_select_{course_idx}_{file_idx}_{course_folder}_{file_name}"
                            if st.button(f"📄 {file_name}", key=unique_key):
                                st.session_state.selected_file_for_ai = {
                                    'course': course_folder,
                                    'file_name': file_name,
                                    'file_path': file_path
                                }
                                st.rerun()
                    else:
                        st.info("No supported files found in this course folder.")
                        
                except Exception as e:
                    st.error(f"Error reading course folder: {e}")
    
    def render_ai_processing_panel_lms_style(self):
        """Render AI processing panel using exact LMS AI Assistant architecture"""
        if not st.session_state.selected_file_for_ai:
            st.info("👈 Please select a file from the course list to generate AI notes.")
            return
        
        selected_file = st.session_state.selected_file_for_ai
        file_key = f"{selected_file['course']}_{selected_file['file_name']}"
        
        # Show selected file info
        st.write(f"**Selected File:** {selected_file['file_name']}")
        st.write(f"**Course:** {selected_file['course']}")
        
        # Notes generation controls (exact LMS AI Assistant style)
        st.markdown("### 📝 Notes Generation Settings")
        
        col_a, col_b = st.columns([3, 1])
        with col_a:
            custom_title = st.text_input("Custom Title (optional)", value=selected_file['file_name'])
        with col_b:
            group_size = st.number_input("Chunks/group", min_value=1, max_value=10, value=3)
        
        # Check if notes are cached
        cached_notes = st.session_state.ai_notes_cache.get(file_key)
        
        if cached_notes:
            st.success("💾 Cached notes found!")
            
            # Option to regenerate
            col1, col2 = st.columns([1, 1])
            with col1:
                if st.button("🔄 Regenerate Notes", type="secondary"):
                    self.generate_ai_notes_lms_style(selected_file, file_key, custom_title, group_size, force_regenerate=True)
            with col2:
                if st.button("🗑️ Clear Cache", type="secondary"):
                    if file_key in st.session_state.ai_notes_cache:
                        del st.session_state.ai_notes_cache[file_key]
                    st.rerun()
            
            # Display cached notes
            st.markdown("### 📝 Generated Notes")
            st.markdown(cached_notes)
            
        else:
            st.info("🤖 No cached notes found for this file.")
            
            # Generate new notes
            if st.button("✨ Generate AI Notes", type="primary"):
                self.generate_ai_notes_lms_style(selected_file, file_key, custom_title, group_size)
    
    def generate_ai_notes_lms_style(self, selected_file, file_key, custom_title, group_size, force_regenerate=False):
        """Generate AI notes using exact LMS AI Assistant architecture"""
        try:
            from utils.notes import iter_generate_notes_from_texts
            from utils.text_processing import chunk_text
            
            with st.spinner("🤖 Generating AI notes using LMS AI Assistant architecture..."):
                file_path = selected_file['file_path']
                
                # Extract text from file
                result = self._extract_text_from_file(Path(file_path))
                
                if result.get('success') and result.get('full_text'):
                    extracted_text = result['full_text']
                    
                    # Use exact LMS AI Assistant chunking
                    chunks = chunk_text(extracted_text, chunk_size=1200, chunk_overlap=200)
                    
                    if not chunks:
                        st.error("❌ No text chunks could be extracted from the file.")
                        return
                    
                    # Generate notes using exact LMS AI Assistant method
                    sections = []
                    placeholder = st.empty()
                    
                    # Use the exact same streaming approach as LMS AI Assistant
                    for idx, sec in enumerate(iter_generate_notes_from_texts(
                        chunks, 
                        title=custom_title, 
                        group_size=int(group_size)
                    ), 1):
                        sections.append(sec)
                        with placeholder.container():
                            for i, s in enumerate(sections, 1):
                                with st.expander(f"Section {i}", expanded=(i == idx)):
                                    st.markdown(s)
                    
                    # Cache the complete notes
                    complete_notes = "\n\n".join(sections)
                    st.session_state.ai_notes_cache[file_key] = complete_notes
                    
                    st.success("✅ Notes generated successfully using LMS AI Assistant architecture!")
                    
                else:
                    error_msg = result.get('error', 'Unknown error')
                    st.error(f"❌ Could not extract text from the file: {error_msg}")
                    
        except Exception as e:
            st.error(f"❌ Error generating notes: {str(e)}")
            st.exception(e)
    
    def render_course_file_selector(self):
        """Render course file selector from downloads folder"""
        downloads_dir = os.path.join(os.path.dirname(__file__), 'downloads')
        
        if not os.path.exists(downloads_dir):
            st.warning("📁 No downloads folder found. Please download some course files first.")
            return
        
        # Get all course folders
        course_folders = [f for f in os.listdir(downloads_dir) 
                         if os.path.isdir(os.path.join(downloads_dir, f))]
        
        if not course_folders:
            st.warning("📁 No course folders found in downloads.")
            return
        
        # Course selection with expandable view
        for course_idx, course_folder in enumerate(course_folders):
            course_path = os.path.join(downloads_dir, course_folder)
            
            with st.expander(f"📚 {course_folder}", expanded=False):
                # Get files in this course folder
                try:
                    files = [f for f in os.listdir(course_path) 
                            if os.path.isfile(os.path.join(course_path, f))]
                    
                    if files:
                        for file_idx, file_name in enumerate(sorted(files)):
                            file_path = os.path.join(course_path, file_name)
                            
                            # Create unique key using course and file indices
                            unique_key = f"select_{course_idx}_{file_idx}_{course_folder}_{file_name}"
                            
                            # File selection button
                            if st.button(f"📄 {file_name}", key=unique_key):
                                st.session_state.selected_file_for_ai = {
                                    'course': course_folder,
                                    'file_name': file_name,
                                    'file_path': file_path
                                }
                                st.rerun()
                    else:
                        st.write("No files found in this course.")
                        
                except Exception as e:
                    st.error(f"Error reading course folder: {e}")
        
        # Show currently selected file
        if st.session_state.selected_file_for_ai:
            selected = st.session_state.selected_file_for_ai
            st.success(f"✅ Selected: {selected['file_name']}")
            st.caption(f"From: {selected['course']}")
    
    def render_ai_processing_panel(self):
        """Render AI processing panel with caching"""
        if not st.session_state.selected_file_for_ai:
            st.info("👈 Please select a file from the course list to generate AI notes.")
            return
        
        selected_file = st.session_state.selected_file_for_ai
        file_key = f"{selected_file['course']}_{selected_file['file_name']}"
        
        # Show selected file info
        st.write(f"**Selected File:** {selected_file['file_name']}")
        st.write(f"**Course:** {selected_file['course']}")
        
        # Check if notes are cached
        cached_notes = st.session_state.ai_notes_cache.get(file_key)
        
        if cached_notes:
            st.success("💾 Cached notes found!")
            
            # Option to regenerate
            col1, col2 = st.columns([1, 1])
            with col1:
                if st.button("🔄 Regenerate Notes", type="secondary"):
                    self.generate_ai_notes(selected_file, file_key, force_regenerate=True)
            with col2:
                if st.button("🗑️ Clear Cache", type="secondary"):
                    if file_key in st.session_state.ai_notes_cache:
                        del st.session_state.ai_notes_cache[file_key]
                    st.rerun()
            
            # Display cached notes with enhanced formatting
            st.markdown("### 📝 Generated Notes")
            
            # Create a beautiful container for the notes
            st.markdown("""
            <div style="
                background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                padding: 20px;
                border-radius: 15px;
                margin: 10px 0;
                box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1);
                border: 1px solid rgba(255, 255, 255, 0.2);
            ">
                <div style="
                    background: white;
                    padding: 20px;
                    border-radius: 10px;
                    box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
                ">
            """, unsafe_allow_html=True)
            
            st.markdown(cached_notes)
            
            st.markdown("""
                </div>
            </div>
            """, unsafe_allow_html=True)
            
        else:
            st.info("🤖 No cached notes found for this file.")
            
            # Generate new notes
            if st.button("✨ Generate AI Notes", type="primary"):
                self.generate_ai_notes(selected_file, file_key)
    
    def generate_ai_notes(self, selected_file, file_key, force_regenerate=False):
        """Generate AI notes for the selected file"""
        try:
            with st.spinner("🤖 Generating AI notes..."):
                file_path = selected_file['file_path']
                
                # Extract text from file
                result = self._extract_text_from_file(Path(file_path))
                
                if result.get('success') and result.get('full_text'):
                    extracted_text = result['full_text']
                    # Generate notes using built-in method
                    notes = self._generate_simple_notes(extracted_text, selected_file['file_name'])
                    
                    # Cache the notes
                    st.session_state.ai_notes_cache[file_key] = notes
                    
                    st.success("✅ Notes generated successfully!")
                    st.rerun()
                else:
                    error_msg = result.get('error', 'Unknown error')
                    
                    # Enhanced logging for debugging
                    logger.error(f"Text extraction failed for {file_path}: {error_msg}")
                    logger.error(f"Full extraction result: {result}")
                    
                    # Categorize and provide specific error messages
                    if 'EOF marker not found' in error_msg or 'PdfReadError' in error_msg:
                        st.error(f"❌ Could not extract text from the file: PDF file is corrupted or has formatting issues")
                        st.warning("💡 This PDF file appears to be corrupted or has formatting issues. Try with a different PDF file.")
                    elif 'permission' in error_msg.lower() or 'denied' in error_msg.lower():
                        st.error(f"❌ Could not extract text from the file: Permission denied")
                        st.warning("💡 Permission denied. Make sure the file is not password-protected or in use by another application.")
                    elif 'not found' in error_msg.lower() or 'FileNotFoundError' in error_msg:
                        st.error(f"❌ Could not extract text from the file: File not found")
                        st.warning("💡 File not found. Please make sure the file exists and try again.")
                    elif 'Unsupported file type' in error_msg:
                        st.error(f"❌ Could not extract text from the file: Unsupported file format")
                        st.warning("💡 This file format is not supported. Try with TXT, PDF, DOCX, or PPTX files.")
                    elif 'password' in error_msg.lower() or 'encrypted' in error_msg.lower():
                        st.error(f"❌ Could not extract text from the file: Password-protected file")
                        st.warning("💡 This file is password-protected. Please remove the password protection and try again.")
                    elif 'empty' in error_msg.lower() or 'no text' in error_msg.lower():
                        st.error(f"❌ Could not extract text from the file: No readable text found")
                        st.warning("💡 This file appears to contain no readable text. It might be an image-based PDF or empty document.")
                    else:
                        # Improved fallback with debugging info
                        st.error(f"❌ Could not extract text from the file: {error_msg}")
                        st.warning("💡 Try with a different file format (TXT, DOCX, or another PDF).")
                        # Show debug info in expander for troubleshooting
                        with st.expander("🔍 Debug Information", expanded=False):
                            st.code(f"Error details: {error_msg}\nFile: {file_path}\nResult: {result}")
        except Exception as e:
            st.error(f"❌ Error generating notes: {str(e)}")
            logger.error(f"Notes generation error: {e}")

    
    def _extract_text_from_file(self, file_path: Path):
        """Simple text extraction from local files"""
        try:
            file_extension = file_path.suffix.lower()
            
            if file_extension == '.txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
                return {'success': True, 'full_text': text, 'filename': file_path.name}
            
            elif file_extension == '.pdf':
                try:
                    import PyPDF2
                    with open(file_path, 'rb') as f:
                        reader = PyPDF2.PdfReader(f)
                        text = ""
                        for page in reader.pages:
                            text += page.extract_text() + "\n"
                    return {'success': True, 'full_text': text, 'filename': file_path.name}
                except ImportError:
                    return {'success': False, 'error': 'PyPDF2 not installed. Install with: pip install PyPDF2'}
                except Exception as e:
                    return {'success': False, 'error': f'PDF extraction failed: {str(e)}'}
            
            elif file_extension in ['.docx']:
                try:
                    import docx
                    doc = docx.Document(file_path)
                    text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                    return {'success': True, 'full_text': text, 'filename': file_path.name}
                except ImportError:
                    return {'success': False, 'error': 'python-docx not installed. Install with: pip install python-docx'}
                except Exception as e:
                    return {'success': False, 'error': f'DOCX extraction failed: {str(e)}'}
            
            elif file_extension in ['.pptx']:
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                    return {'success': True, 'full_text': text, 'filename': file_path.name}
                except ImportError:
                    return {'success': False, 'error': 'python-pptx not installed. Install with: pip install python-pptx'}
                except Exception as e:
                    return {'success': False, 'error': f'PPTX extraction failed: {str(e)}'}
            
            else:
                return {'success': False, 'error': f'Unsupported file type: {file_extension}'}
                
        except Exception as e:
            return {'success': False, 'error': f'File extraction failed: {str(e)}'}

    def process_files_for_notes(self, files: List[Path], title: str, embed_model: str, top_k: int, chunk_size: int, group_size: int):
        """Process local files and generate AI notes"""
        try:
            with st.spinner("Processing local files and generating notes..."):
                # Extract text from local files
                extraction_results = []
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                # Filter for existing local files
                existing_files = [f for f in files if f.exists()]
                if not existing_files:
                    st.error("❌ No local files found. Please download files first using the Download section.")
                    st.info("💡 **Tip:** Use the 'Download Files from JSON Exports' section to get files locally before processing.")
                    return
                
                if len(existing_files) < len(files):
                    missing_files = [f for f in files if not f.exists()]
                    st.warning(f"⚠️ {len(missing_files)} files not found locally. Processing {len(existing_files)} available files.")
                
                for i, file_path in enumerate(existing_files):
                    status_text.text(f"Extracting text from local file: {file_path.name}")
                    
                    try:
                        # Check file size
                        file_size = file_path.stat().st_size
                        if file_size == 0:
                            st.warning(f"⚠️ Skipping {file_path.name}: Empty file")
                            continue
                        
                        result = self._extract_text_from_file(file_path)
                        if result.get('success'):
                            # Add local file metadata
                            result['local_file_path'] = str(file_path)
                            result['file_size'] = file_size
                            extraction_results.append(result)
                        else:
                            error_msg = result.get('error', 'Unknown error')
                            
                            # Enhanced error categorization for local files
                            if 'EOF marker not found' in error_msg or 'PdfReadError' in error_msg:
                                st.warning(f"⚠️ Skipping {file_path.name}: PDF file is corrupted or has formatting issues")
                            elif 'permission' in error_msg.lower() or 'denied' in error_msg.lower():
                                st.warning(f"⚠️ Skipping {file_path.name}: Permission denied (file may be password-protected)")
                            elif 'Unsupported file type' in error_msg:
                                st.warning(f"⚠️ Skipping {file_path.name}: Unsupported file format")
                            elif 'password' in error_msg.lower() or 'encrypted' in error_msg.lower():
                                st.warning(f"⚠️ Skipping {file_path.name}: Password-protected file")
                            elif 'empty' in error_msg.lower() or 'no text' in error_msg.lower():
                                st.warning(f"⚠️ Skipping {file_path.name}: No readable text found")
                            else:
                                st.warning(f"⚠️ Failed to extract from local file {file_path.name}: {error_msg}")
                    except Exception as e:
                        st.warning(f"⚠️ Failed to extract from local file {file_path.name}: {e}")
                    
                    progress_bar.progress((i + 1) / len(existing_files))
                
                if not extraction_results:
                    st.error("❌ No local files were successfully processed")
                    st.info("💡 **Tip:** Make sure files are downloaded locally and are not password-protected or corrupted.")
                    return
                
                # Show processing summary
                st.success(f"✅ Successfully processed {len(extraction_results)} local files")
                
                # Display file processing summary
                with st.expander("📊 Local File Processing Summary", expanded=True):
                    total_size = sum(result.get('file_size', 0) for result in extraction_results)
                    st.write(f"**Files Processed:** {len(extraction_results)}")
                    st.write(f"**Total Size:** {total_size / (1024*1024):.2f} MB")
                    st.write(f"**Model:** {embed_model}")
                    st.write(f"**Chunk Size:** {chunk_size}")
                    st.write(f"**Group Size:** {group_size}")
                    
                    # Show processed files
                    st.write("**Processed Local Files:**")
                    for result in extraction_results:
                        file_size_mb = result.get('file_size', 0) / (1024 * 1024)
                        st.write(f"• {Path(result.get('local_file_path', '')).name} ({file_size_mb:.2f} MB)")
                
                status_text.text("Generating AI notes from local files...")
                
                # Generate notes from extracted text
                all_texts = []
                for result in extraction_results:
                    if 'full_text' in result:
                        all_texts.append(result['full_text'])
                
                if all_texts:
                    # Simple note generation without external dependencies
                    st.subheader("📝 Generated Notes from Local Files")
                    
                    # Combine all text
                    combined_text = "\n\n".join(all_texts)
                    
                    # Simple text processing and note generation
                    notes = self._generate_simple_notes(combined_text, title)
                    
                    # Display notes with enhanced formatting
                    st.markdown("""
                    <div style="
                        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                        padding: 20px;
                        border-radius: 15px;
                        margin: 10px 0;
                        box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1);
                        border: 1px solid rgba(255, 255, 255, 0.2);
                    ">
                        <div style="
                            background: white;
                            padding: 20px;
                            border-radius: 10px;
                            box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
                        ">
                    """, unsafe_allow_html=True)
                    
                    st.markdown(notes)
                    
                    st.markdown("""
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    # Save notes to file
                    notes_dir = self.download_dir / "ai_notes"
                    notes_dir.mkdir(exist_ok=True)
                    
                    # Save notes
                    safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).rstrip()
                    safe_title = safe_title.replace(' ', '_')
                    notes_filename = f"{safe_title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.md"
                    notes_filepath = notes_dir / notes_filename
                    
                    with open(notes_filepath, 'w', encoding='utf-8') as f:
                        f.write(f"# {title}\n\n")
                        f.write(f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")
                        f.write(f"Source: Local files\n")
                        f.write(f"Files processed: {', '.join([Path(result.get('local_file_path', '')).name for result in extraction_results])}\n\n")
                        f.write("---\n\n")
                        f.write(notes)
                    
                    st.success(f"✅ Notes saved to: {notes_filepath}")
                    
                    # Provide download button
                    with open(notes_filepath, 'rb') as f:
                        st.download_button(
                            label="📥 Download Notes (Markdown)",
                            data=f.read(),
                            file_name=notes_filename,
                            mime="text/markdown"
                        )
                
                status_text.text("Notes generation completed!")
                
        except Exception as e:
            st.error(f"❌ Error generating notes: {str(e)}")
            logger.error(f"Notes generation error: {e}")
    
    def check_course_files_availability(self, course_id: int) -> Dict[str, Any]:
        """Check if course files are available in downloads directory"""
        course_pattern = f"course_{course_id}_*"
        course_dirs = list(self.download_dir.glob(course_pattern))
        
        if not course_dirs:
            return {
                'available': False,
                'directories': [],
                'file_count': 0
            }
        
        total_files = 0
        for course_dir in course_dirs:
            # Count supported files
            supported_extensions = ['.pdf', '.docx', '.doc', '.txt', '.pptx', '.ppsx']
            for ext in supported_extensions:
                total_files += len(list(course_dir.rglob(f'*{ext}')))
        
        return {
            'available': True,
            'directories': course_dirs,
            'file_count': total_files
        }
    
    def check_downloads_folder_availability(self) -> Dict[str, Any]:
        """Check if files are available in the main downloads folder"""
        downloads_path = Path("/Users/praveenbhandari/sf-vibe/downloads")
        
        if not downloads_path.exists():
            return {
                'available': False,
                'directories': [],
                'file_count': 0
            }
        
        # Find all course directories in downloads
        course_dirs = [d for d in downloads_path.iterdir() if d.is_dir()]
        
        if not course_dirs:
            return {
                'available': False,
                'directories': [],
                'file_count': 0
            }
        
        total_files = 0
        for course_dir in course_dirs:
            # Count supported files
            supported_extensions = ['.pdf', '.docx', '.doc', '.txt', '.pptx', '.ppsx']
            for ext in supported_extensions:
                total_files += len(list(course_dir.rglob(f'*{ext}')))
        
        return {
            'available': total_files > 0,
            'directories': course_dirs,
            'file_count': total_files
        }
    
    def render_enhanced_download_section(self):
        """Enhanced download section with file availability checking"""
        st.header("📥 Course File Downloads")
        
        if not st.session_state.get('selected_course'):
            st.warning("⚠️ Please select a course first.")
            return
            
        course_id = st.session_state.selected_course['id']
        course_name = st.session_state.selected_course['name']
        
        # Check file availability in downloads folder
        file_status = self.check_downloads_folder_availability()
        
        # Status display
        col1, col2, col3 = st.columns(3)
        with col1:
            if file_status['available']:
                st.success(f"✅ Files Available")
            else:
                st.warning("⚠️ No Files Downloaded")
        
        with col2:
            st.metric("Downloaded Directories", len(file_status['directories']))
        
        with col3:
            st.metric("Total Files", file_status['file_count'])
        
        # Show existing downloads if available
        if file_status['available']:
            st.subheader("📁 Existing Downloads")
            
            for course_dir in file_status['directories']:
                with st.expander(f"📂 {course_dir.name}", expanded=False):
                    # Show directory contents
                    supported_extensions = ['.pdf', '.docx', '.doc', '.txt', '.pptx', '.ppsx']
                    files_found = []
                    
                    for ext in supported_extensions:
                        files_found.extend(course_dir.rglob(f'*{ext}'))
                    
                    if files_found:
                        st.markdown(f"**📊 Found {len(files_found)} supported files:**")
                        
                        # Create a grid layout for files
                        cols = st.columns(2)
                        for idx, file_path in enumerate(files_found[:10]):  # Show first 10 files
                            relative_path = file_path.relative_to(course_dir)
                            col = cols[idx % 2]
                            
                            with col:
                                # Get file extension for appropriate icon
                                file_ext = relative_path.suffix.lower().replace('.', '') if relative_path.suffix else ''
                                file_icons = {
                                    'pdf': '📕',
                                    'doc': '📄',
                                    'docx': '📄',
                                    'txt': '📝',
                                    'ppt': '📊',
                                    'pptx': '📊',
                                    'xls': '📈',
                                    'xlsx': '📈',
                                    'jpg': '🖼️',
                                    'jpeg': '🖼️',
                                    'png': '🖼️',
                                    'gif': '🖼️',
                                    'mp4': '🎥',
                                    'mp3': '🎵',
                                    'zip': '📦'
                                }
                                icon = file_icons.get(file_ext, '📄')
                                
                                # Create a styled file entry
                                st.markdown(f"""
                                <div style="
                                    background-color: #f8f9fa;
                                    padding: 8px 12px;
                                    border-radius: 6px;
                                    margin: 3px 0;
                                    border-left: 3px solid #007bff;
                                    font-size: 14px;
                                ">
                                    {icon} <strong>{relative_path.name}</strong><br>
                                    <small style="color: #6c757d;">{relative_path.parent if relative_path.parent != Path('.') else 'Root'}</small>
                                </div>
                                """, unsafe_allow_html=True)
                        
                        # Show "and X more" if there are more files with better styling
                        if len(files_found) > 10:
                            st.markdown(f"""
                            <div style="
                                text-align: center;
                                padding: 12px;
                                background-color: #e3f2fd;
                                border-radius: 8px;
                                margin: 15px 0;
                                color: #1976d2;
                                font-weight: bold;
                                border: 1px dashed #1976d2;
                            ">
                                📁 ... and {len(files_found) - 10} more files
                            </div>
                            """, unsafe_allow_html=True)
                        
                        # Quick action buttons
                        col1, col2 = st.columns(2)
                        with col1:
                            if st.button(f"🤖 Generate AI Notes", key=f"ai_notes_{course_dir.name}"):
                                st.info("Switching to AI Notes tab...")
                                st.rerun()
                        
                        with col2:
                            if st.button(f"📂 Open Folder", key=f"open_{course_dir.name}"):
                                st.info(f"Folder location: {course_dir}")
                    else:
                        st.write("No supported files found in this directory.")
            
            st.divider()
        
        # Download new files section
        st.subheader("📥 Download Course Files")
        
        if not file_status['available']:
            st.info("💡 **Recommended**: Download course files to enable AI note generation and offline access.")
        
        # Check if canvas_client is available for new downloads
        if not st.session_state.get('canvas_client'):
            st.warning("⚠️ Please log in first to enable new file downloads.")
            if file_status['available']:
                st.info("✅ You can still use the existing files above for AI note generation.")
            return
        
        # Initialize downloader
        if 'downloader' not in st.session_state:
            # Get API token from canvas client if available
            api_token = None
            base_url = None
            if hasattr(st.session_state, 'canvas_client') and st.session_state.canvas_client:
                api_token = getattr(st.session_state.canvas_client, '_Canvas__requester', {}).get('access_token')
                base_url = getattr(st.session_state.canvas_client, 'base_url', None)
            
            st.session_state.downloader = CanvasFileDownloader(
                api_token=api_token,
                base_url=base_url,
                output_dir=str(self.download_dir),
                dry_run=False
            )
        
        downloader = st.session_state.downloader
        
        # Download options
        col1, col2 = st.columns(2)
        
        with col1:
            st.write("**Download Options:**")
            download_files = st.checkbox("📄 Course Files", value=True)
            download_modules = st.checkbox("📚 Module Content", value=True)
            download_assignments = st.checkbox("📝 Assignment Files", value=False)
        
        with col2:
            st.write("**Settings:**")
            dry_run = st.checkbox("🔍 Dry Run (Preview Only)", value=False)
            organize_by_type = st.checkbox("📁 Organize by File Type", value=True)
            max_file_size = st.number_input("Max File Size (MB)", min_value=1, max_value=500, value=100)
        
        # Custom download directory
        custom_dir = st.text_input(
            "Custom Download Directory (optional):",
            placeholder=str(self.download_dir)
        )
        
        if custom_dir:
            download_dir = Path(custom_dir)
        else:
            download_dir = self.download_dir
        
        # Download button
        if st.button("🚀 Start Download", type="primary"):
            if not any([download_files, download_modules, download_assignments]):
                st.error("❌ Please select at least one download option.")
                return
            
            # Update downloader settings
            downloader.download_dir = download_dir
            downloader.max_file_size_mb = max_file_size
            
            # Start download process
            with st.spinner("Preparing download..."):
                try:
                    # Create progress containers
                    progress_container = st.container()
                    log_container = st.container()
                    
                    with progress_container:
                        progress_bar = st.progress(0)
                        status_text = st.empty()
                        stats_container = st.empty()
                    
                    # Download course content
                    course = st.session_state.selected_course
                    
                    if dry_run:
                        status_text.text("🔍 Running dry run - no files will be downloaded")
                        st.info("Dry run mode: This will show what would be downloaded without actually downloading files.")
                    
                    # Process downloads
                    total_downloaded = 0
                    total_failed = 0
                    downloaded_files = []
                    failed_downloads = []
                    
                    if download_files:
                        status_text.text("📄 Processing course files...")
                        try:
                            result = downloader.download_course_files(
                                course_id=course['id'],
                                course_name=course['name'],
                                dry_run=dry_run,
                                organize_by_type=organize_by_type
                            )
                            total_downloaded += result.get('downloaded', 0)
                            total_failed += result.get('failed', 0)
                            downloaded_files.extend(result.get('files', []))
                            failed_downloads.extend(result.get('failed_files', []))
                        except Exception as e:
                            st.error(f"Error downloading course files: {e}")
                        
                        progress_bar.progress(0.33)
                    
                    if download_modules:
                        status_text.text("📚 Processing module content...")
                        try:
                            # Get modules
                            if not st.session_state.get('canvas_client'):
                                st.error("Canvas client not available. Please log in again.")
                                return
                            modules = list(st.session_state.canvas_client.get_course(course['id']).get_modules())
                            
                            for i, module in enumerate(modules):
                                module_items = list(module.get_module_items())
                                for item in module_items:
                                    if hasattr(item, 'url') and item.url:
                                        # Process module item
                                        try:
                                            downloader.process_module_item(
                                                item,
                                                course['id'],
                                                course['name'],
                                                module.name,
                                                dry_run=dry_run
                                            )
                                            total_downloaded += 1
                                        except Exception as e:
                                            total_failed += 1
                                            failed_downloads.append(f"Module item: {item.title} - {str(e)}")
                                
                                # Update progress
                                progress = 0.33 + (0.33 * (i + 1) / len(modules))
                                progress_bar.progress(min(progress, 0.66))
                        except Exception as e:
                            st.error(f"Error downloading module content: {e}")
                        
                        progress_bar.progress(0.66)
                    
                    if download_assignments:
                        status_text.text("📝 Processing assignment files...")
                        try:
                            # This would require additional implementation
                            # for assignment-specific file downloads
                            st.info("Assignment file downloads not yet implemented")
                        except Exception as e:
                            st.error(f"Error downloading assignment files: {e}")
                        
                        progress_bar.progress(1.0)
                    
                    # Final status
                    status_text.text("✅ Download process completed!")
                    
                    # Show results
                    with stats_container:
                        col1, col2, col3 = st.columns(3)
                        with col1:
                            st.metric("✅ Downloaded", total_downloaded)
                        with col2:
                            st.metric("❌ Failed", total_failed)
                        with col3:
                            success_rate = (total_downloaded / (total_downloaded + total_failed) * 100) if (total_downloaded + total_failed) > 0 else 0
                            st.metric("Success Rate", f"{success_rate:.1f}%")
                    
                    # Show failed downloads if any
                    if failed_downloads:
                        with st.expander(f"❌ Failed Downloads ({len(failed_downloads)})", expanded=False):
                            for failure in failed_downloads[:20]:  # Show first 20 failures
                                st.write(f"• {failure}")
                            if len(failed_downloads) > 20:
                                st.write(f"... and {len(failed_downloads) - 20} more failures")
                    
                    # Next steps
                    if total_downloaded > 0 and not dry_run:
                        st.success("🎉 Files downloaded successfully!")
                        
                        col1, col2 = st.columns(2)
                        with col1:
                            if st.button("🤖 Generate AI Notes Now"):
                                st.info("Switching to AI Notes tab...")
                                st.rerun()
                        
                        with col2:
                            if st.button("📂 View Downloads Folder"):
                                st.info(f"Downloads saved to: {download_dir}")
                                # Show updated file status
                                updated_status = self.check_course_files_availability(course_id)
                                st.write(f"Total files now available: {updated_status['file_count']}")
                    
                except Exception as e:
                    st.error(f"❌ Download failed: {str(e)}")
                    logger.error(f"Download error: {e}")
        
        # Help section
        with st.expander("ℹ️ Help & Tips", expanded=False):
            st.markdown("""
            ### Download Tips:
            
            **File Types Supported for AI Notes:**
            - 📄 PDF documents
            - 📝 Word documents (.docx, .doc)
            - 📄 Text files (.txt)
            
            **Workflow Recommendations:**
            1. **First Time**: Download course files using this tab
            2. **Generate Notes**: Use the AI Notes tab to process downloaded files
            3. **Organize**: Files are automatically organized by course and type
            
            **Troubleshooting:**
            - Use "Dry Run" to preview what will be downloaded
            - Check Canvas permissions if downloads fail
            - Large files may take longer to process
            
            **File Organization:**
            - Files are saved to: `downloads/course_{id}_{name}/`
            - AI notes are saved to: `downloads/ai_notes/`
            """)
    
    def render_rag_demo_section(self):
        """Render enhanced RAG Demo section with memory and learning features"""
        st.header("🤖 Enhanced RAG Demo - AI Learning Assistant")
        st.markdown("Upload documents or use local files to create an AI-powered learning experience with memory and personalized tutoring.")
        
        # Initialize memory store
        if 'memory_store' not in st.session_state:
            st.session_state.memory_store = self.MemoryStore()
        
        # Memory profile selection
        st.sidebar.subheader("Profile & Memory")
        profile_id = st.sidebar.text_input("Profile ID", value="default")
        memory_short_window = st.sidebar.slider("Short-term window", 4, 12, 6)
        
        # Check for local files in downloads directory
        downloads_dir = Path("downloads")
        local_files = []
        
        if downloads_dir.exists():
            supported_extensions = ['.pdf', '.docx', '.doc', '.txt', '.pptx', '.ppsx']
            for ext in supported_extensions:
                local_files.extend(downloads_dir.rglob(f'*{ext}'))
        
        if not local_files and not st.session_state.uploaded_files:
            st.info("📁 **No Local Files Found**")
            st.markdown("""
            **To use the RAG demo, you need to:**
            1. Download files locally using the Download section
            2. Or upload files manually using the file uploader below
            
            This approach works around Canvas API permission restrictions by processing files locally.
            """)
            
            # File uploader as fallback
            uploaded_files = st.file_uploader(
                "Upload files for RAG processing",
                type=['pdf', 'docx', 'doc', 'txt', 'pptx', 'ppsx'],
                accept_multiple_files=True,
                help="Upload files if you don't have any in the downloads folder"
            )
            
            if uploaded_files:
                st.session_state.uploaded_files = []
                for uploaded_file in uploaded_files:
                    # Save uploaded file temporarily
                    temp_path = Path("temp") / uploaded_file.name
                    temp_path.parent.mkdir(exist_ok=True)
                    with open(temp_path, "wb") as f:
                        f.write(uploaded_file.getbuffer())
                    
                    # Extract text
                    result = self._extract_text_from_file(temp_path)
                    if result.get('success'):
                        st.session_state.uploaded_files.append({
                            'name': uploaded_file.name,
                            'text': result.get('full_text', ''),
                            'type': 'uploaded_document',
                            'file_path': str(temp_path)
                        })
                    temp_path.unlink()  # Clean up temp file
                
                if st.session_state.uploaded_files:
                    st.success(f"✅ Processed {len(st.session_state.uploaded_files)} uploaded files!")
        
        # Show available files
        all_files = []
        if local_files:
            for file_path in local_files[:10]:  # Limit to first 10 files
                result = self._extract_text_from_file(file_path)
                if result.get('success'):
                    all_files.append({
                        'name': file_path.name,
                        'text': result.get('full_text', ''),
                        'type': 'local_document',
                        'file_path': str(file_path)
                    })
        
        if st.session_state.uploaded_files:
            all_files.extend(st.session_state.uploaded_files)
        
        if not all_files:
            return
        
        # Create tabs for different RAG features
        rag_tab1, rag_tab2, rag_tab3, rag_tab4 = st.tabs(["📄 Document Browser", "💬 RAG Chatbot", "🎓 Learning Mode", "🔗 Resources"])
        
        with rag_tab1:
            self.render_document_browser(all_files)
        
        with rag_tab2:
            self.render_rag_chatbot(all_files, profile_id, memory_short_window)
        
        with rag_tab3:
            self.render_learning_mode(all_files, profile_id, memory_short_window)
        
        with rag_tab4:
            self.render_resource_recommendations(all_files)
    
    def render_document_browser(self, all_files):
        """Render document browser using the course file selector"""
        st.subheader("📄 Document Browser")
        
        # Create the beautiful header display
        st.markdown("""
        <div style="
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            padding: 30px;
            border-radius: 15px;
            margin: 20px 0;
            box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1);
            border: 1px solid rgba(255, 255, 255, 0.2);
            text-align: center;
        ">
            <h2 style="color: white; margin: 0; font-size: 28px;">🤖 AI-Powered Note Generation</h2>
            <p style="color: white; margin: 10px 0 0 0; font-size: 18px; opacity: 0.9;">📁 Course Files</p>
        </div>
        """, unsafe_allow_html=True)
        
        # Use the existing course file selector
        with st.container():
            st.subheader("📁 Course Files")
            self.render_course_file_selector()
        
        # Show content preview and actions for selected file
        if hasattr(st.session_state, 'selected_file_for_ai') and st.session_state.selected_file_for_ai:
            selected = st.session_state.selected_file_for_ai
            file_path = selected['file_path']
            
            # Extract text from the selected file
            result = self._extract_text_from_file(Path(file_path))
            
            if result.get('success'):
                st.markdown("---")
                st.markdown("### 📄 Content Preview")
                
                # Create a beautiful container for the content preview
                st.markdown("""
                <div style="
                    background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                    padding: 20px;
                    border-radius: 15px;
                    margin: 10px 0;
                    box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1);
                    border: 1px solid rgba(255, 255, 255, 0.2);
                ">
                    <div style="
                        background: white;
                        padding: 20px;
                        border-radius: 10px;
                        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
                        max-height: 300px;
                        overflow-y: auto;
                    ">
                """, unsafe_allow_html=True)
                
                # Format the content preview
                preview_text = result.get('full_text', '')[:1000] + "..." if len(result.get('full_text', '')) > 1000 else result.get('full_text', '')
                formatted_preview = self._format_content_preview(preview_text, selected['file_name'])
                st.markdown(formatted_preview)
                
                st.markdown("""
                    </div>
                </div>
                """, unsafe_allow_html=True)
                
                # Action buttons for selected file
                col1, col2, col3 = st.columns([1, 1, 1])
                
                with col1:
                    if st.button("📝 Generate AI Notes", key=f"notes_{selected['file_name']}_{selected['course']}", type="primary"):
                        # Create file info for note generation
                        file_info = {
                            'name': selected['file_name'],
                            'text': result.get('full_text', ''),
                            'type': 'local_document',
                            'file_path': file_path
                        }
                        st.session_state.selected_file_for_notes = file_info
                        st.success("File selected for AI note generation!")
                
                with col2:
                    if st.button("👁️ Quick Preview", key=f"preview_{selected['file_name']}_{selected['course']}"):
                        st.info(f"Preview: {selected['file_name']}")
                
                with col3:
                    if st.button("📊 File Info", key=f"info_{selected['file_name']}_{selected['course']}"):
                        file_size = len(result.get('full_text', '')) / 1024  # KB
                        st.info(f"File: {selected['file_name']}\nSize: {file_size:.1f} KB\nCourse: {selected['course']}")
                
                # AI Processing section
                st.markdown("""
                <div style="
                    background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%);
                    padding: 20px;
                    border-radius: 10px;
                    margin: 15px 0;
                    box-shadow: 0 4px 15px rgba(0, 0, 0, 0.1);
                    text-align: center;
                ">
                    <h4 style="color: white; margin: 0; font-size: 18px;">🤖 AI Processing</h4>
                    <p style="color: white; margin: 10px 0 0 0; opacity: 0.9;">Ready to generate AI notes from the selected file!</p>
                </div>
                """, unsafe_allow_html=True)
            else:
                st.error(f"❌ Error reading file: {result.get('error', 'Unknown error')}")
    
    def render_rag_chatbot(self, all_files, profile_id, memory_short_window):
        """Render RAG Chatbot with Memory"""
        st.subheader("💬 RAG Chatbot with Memory")
        
        # Initialize chat history
        if "chat_messages" not in st.session_state:
            st.session_state.chat_messages = []
        if "chat_cooldown" not in st.session_state:
            st.session_state.chat_cooldown = 0.0

        # Chat container
        chat_container = st.container()
        
        # Display chat history in scrollable container
        with chat_container:
            for msg in st.session_state.chat_messages:
                if msg.get("role") == "user":
                    st.markdown(f"**You:** {msg.get('content','')}")
                else:
                    st.markdown(f"**Assistant:** {msg.get('content','')}")

        # Chat input at bottom
        st.markdown("---")
        col1, col2 = st.columns([4, 1])
        with col1:
            user_q = st.text_input("Ask a question about your documents...", key="chat_input", placeholder="Type your question here...")
        with col2:
            send_button = st.button("Send", key="chat_send", type="primary")
        
        if send_button and user_q:
            # Simple cooldown to avoid hammering
            now = time.time()
            if now - st.session_state.chat_cooldown < 3.0:
                st.warning("Please wait a moment before sending another message.")
                st.stop()
            
            with st.spinner("Retrieving and answering..."):
                # Simple document search
                relevant_docs = []
                for doc in all_files:
                    if user_q.lower() in doc['text'].lower():
                        relevant_docs.append({
                            "source": doc['name'],
                            "chunk_index": 0,
                            "text": doc['text'][:1000] + "..." if len(doc['text']) > 1000 else doc['text']
                        })
                
                # Memory contexts
                mem_ctx = st.session_state.memory_store.memory_contexts(st.session_state.chat_messages, profile_id=profile_id, short_window=memory_short_window)
                all_ctx = relevant_docs + mem_ctx
                
                if not all_ctx:
                    answer = "No relevant information found in the documents. Please try a different question or upload more documents."
                else:
                    try:
                        answer = st.session_state.memory_store._answer_with_context(user_q, all_ctx)
                    except Exception as e:
                        answer = f"Error generating response: {str(e)}"
                
                st.session_state.chat_messages.append({"role": "user", "content": user_q})
                st.session_state.chat_messages.append({"role": "assistant", "content": answer})
                
                # Try to update long-term memory occasionally
                _ = st.session_state.memory_store.summarize_and_store_long_term(st.session_state.chat_messages, profile_id=profile_id)
                st.session_state.chat_cooldown = time.time()
                st.rerun()
    
    def render_learning_mode(self, all_files, profile_id, memory_short_window):
        """Render Learning Mode with AI Tutor"""
        st.subheader("🎓 Learning Mode - AI Tutor")
        
        # Initialize learning chat
        if "learning_messages" not in st.session_state:
            st.session_state.learning_messages = []
        if "learning_cooldown" not in st.session_state:
            st.session_state.learning_cooldown = 0.0
        if "learning_initialized" not in st.session_state:
            st.session_state.learning_initialized = False

        # Initialize tutor with topics from documents
        if not st.session_state.learning_initialized:
            # Extract topics from document content
            all_texts = [doc['text'] for doc in all_files]
            topics = self.extract_topics_from_notes(all_texts)
            if not topics:
                topics = [f"Topic {i+1}" for i in range(min(5, len(all_files)))]
            
            st.session_state.learning_topics = topics[:10]
            st.session_state.learning_initialized = True
            
            # Welcome message
            welcome_msg = f"Hello! I'm your AI learning tutor. I can help you learn about these topics: {', '.join(topics[:5])}. What would you like to explore first?"
            st.session_state.learning_messages.append({"role": "assistant", "content": welcome_msg})

        # Chat container
        learning_container = st.container()
        
        # Display chat history
        with learning_container:
            for msg in st.session_state.learning_messages:
                if msg.get("role") == "user":
                    st.markdown(f"**You:** {msg.get('content','')}")
                else:
                    st.markdown(f"**Tutor:** {msg.get('content','')}")

        # Chat input at bottom
        st.markdown("---")
        col1, col2 = st.columns([4, 1])
        with col1:
            user_input = st.text_input("Ask your tutor anything...", key="learning_input", placeholder="What would you like to learn about?")
        with col2:
            send_button = st.button("Send", key="learning_send", type="primary")
        
        if send_button and user_input:
            # Cooldown check
            now = time.time()
            if now - st.session_state.learning_cooldown < 3.0:
                st.warning("Please wait a moment before sending another message.")
                st.stop()
            
            with st.spinner("Tutor is thinking..."):
                # Build context from documents
                relevant_docs = []
                for doc in all_files:
                    if user_input.lower() in doc['text'].lower():
                        relevant_docs.append({
                            "source": doc['name'],
                            "chunk_index": 0,
                            "text": doc['text'][:1000] + "..." if len(doc['text']) > 1000 else doc['text']
                        })
                
                # Memory context
                mem_ctx = st.session_state.memory_store.memory_contexts(st.session_state.learning_messages, profile_id=profile_id, short_window=memory_short_window)
                
                # Get recommended resources
                relevant_topic = user_input
                for topic in st.session_state.learning_topics:
                    if topic.lower() in user_input.lower() or user_input.lower() in topic.lower():
                        relevant_topic = topic
                        break
                
                vlinks = self.recommend_youtube_ddg(relevant_topic)
                alinks = self.recommend_articles_ddg(relevant_topic)
                res_text = f"Recommended Videos for '{relevant_topic}':\n" + "\n".join([f"- {l['title']}: {l['url']}" for l in vlinks[:2]]) + \
                              f"\nRecommended Articles for '{relevant_topic}':\n" + "\n".join([f"- {a['title']}: {a['url']}" for a in alinks[:2]])
                res_ctx = [{"source": "resources", "chunk_index": 0, "text": res_text}]
                
                all_ctx = relevant_docs + mem_ctx + res_ctx
                
                # Tutor prompt
                prompt = f"You are an AI learning tutor. Help the student learn by explaining concepts, asking questions, suggesting resources, and giving small assignments. Be encouraging and educational. Student said: {user_input}"
                
                try:
                    reply = st.session_state.memory_store._answer_with_context(prompt, all_ctx)
                except Exception as e:
                    reply = f"I'm having trouble connecting right now. Please try again in a moment. Error: {e}"
                
                st.session_state.learning_messages.append({"role": "user", "content": user_input})
                st.session_state.learning_messages.append({"role": "assistant", "content": reply})
                
                # Update long-term memory
                _ = st.session_state.memory_store.summarize_and_store_long_term(st.session_state.learning_messages, profile_id=profile_id)
                st.session_state.learning_cooldown = time.time()
                st.rerun()
    
    def render_resource_recommendations(self, all_files):
        """Render Resource Recommendations"""
        st.subheader("🔗 Recommended Resources")
        
        # Build recommendations per detected topic from documents
        all_texts = [doc['text'] for doc in all_files]
        topics = self.extract_topics_from_notes(all_texts)
        if not topics:
            st.info("Could not detect headings; using first few sections as topics.")
            topics = [f"Topic {i+1}" for i in range(min(5, len(all_files)))]
        
        for topic in topics[:10]:
            with st.expander(topic, expanded=False):
                st.markdown("**Related Videos (via DuckDuckGo)**")
                vlinks = self.recommend_youtube_ddg(topic)
                if vlinks:
                    for l in vlinks:
                        st.write(f"- [{l['title']}]({l['url']})")
                else:
                    st.write("No videos found for this topic.")
                
                st.markdown("**Related Articles (via DuckDuckGo)**")
                alinks = self.recommend_articles_ddg(topic)
                if alinks:
                    for a in alinks:
                        st.write(f"- [{a['title']}]({a['url']})")
                else:
                    st.write("No articles found for this topic.")
    
    def render_document_ingestion(self):
        """Render document ingestion interface"""
        st.subheader("📄 Document Ingestion")
        
        # Course files browser (replacing file upload)
        st.markdown("### 📁 Course Files")
        
        downloads_dir = os.path.join(os.path.dirname(__file__), 'downloads')
        
        if not os.path.exists(downloads_dir):
            st.warning("📁 No downloads folder found. Please download some course files first.")
            st.info("💡 Go to the 'Download Files' tab to download course materials first.")
            return
        
        # Get all course folders
        course_folders = [f for f in os.listdir(downloads_dir) 
                         if os.path.isdir(os.path.join(downloads_dir, f))]
        
        if not course_folders:
            st.warning("📁 No course folders found in downloads.")
            st.info("💡 Go to the 'Download Files' tab to download course materials first.")
            return
        
        # Course selection with expandable view
        selected_files = []
        
        for course_folder in course_folders:
            course_path = os.path.join(downloads_dir, course_folder)
            
            with st.expander(f"📚 {course_folder}", expanded=False):
                # Get files in this course folder
                try:
                    files = [f for f in os.listdir(course_path) 
                            if os.path.isfile(os.path.join(course_path, f)) and 
                            f.lower().endswith(('.pdf', '.docx', '.txt'))]
                    
                    if files:
                        st.markdown(f"**Found {len(files)} supported files:**")
                        
                        for file_idx, file_name in enumerate(sorted(files)):
                            file_path = os.path.join(course_path, file_name)
                            
                            # File selection checkbox with unique key
                            unique_key = f"rag_select_{course_folders.index(course_folder)}_{file_idx}_{course_folder}_{file_name}"
                            if st.checkbox(f"📄 {file_name}", key=unique_key):
                                selected_files.append({
                                    'name': file_name,
                                    'path': file_path,
                                    'course': course_folder
                                })
                    else:
                        st.write("No supported files (PDF, DOCX, TXT) found in this course.")
                        
                except Exception as e:
                    st.error(f"Error reading course folder: {e}")
        
        # Store selected files in session state
        if selected_files:
            st.session_state.rag_selected_files = selected_files
            st.success(f"✅ Selected {len(selected_files)} files for processing")
            
            # Show selected files summary
            with st.expander("📋 Selected Files Summary", expanded=True):
                for file_info in selected_files:
                    st.write(f"• **{file_info['name']}** from {file_info['course']}")
        else:
            st.session_state.rag_selected_files = []
        
        # YouTube URL input
        youtube_url = st.text_input(
            "YouTube URL (optional)",
            placeholder="https://www.youtube.com/watch?v=...",
            help="Add educational videos to your knowledge base"
        )
        
        if st.button("🚀 Process Documents", type="primary"):
            selected_files = st.session_state.get('rag_selected_files', [])
            
            if selected_files or youtube_url:
                with st.spinner("Processing documents..."):
                    try:
                        # Process selected course files
                        if selected_files:
                            for file_info in selected_files:
                                # Extract text from course file using built-in method
                                result = self._extract_text_from_file(Path(file_info['path']))
                                
                                if result.get('success') and result.get('full_text'):
                                    text = result['full_text']
                                    st.session_state.uploaded_files.append({
                                        'name': f"{file_info['name']} ({file_info['course']})",
                                        'text': text,
                                        'type': 'course_document'
                                    })
                                else:
                                    st.warning(f"⚠️ Could not extract text from {file_info['name']}: {result.get('error', 'Unknown error')}")
                        
                        # Process YouTube URL
                        if youtube_url:
                            try:
                                transcript = get_youtube_transcript(youtube_url)
                                if transcript:
                                    st.session_state.uploaded_files.append({
                                        'name': f"YouTube: {youtube_url}",
                                        'text': transcript,
                                        'type': 'video'
                                    })
                            except Exception as e:
                                st.error(f"Error processing YouTube URL: {e}")
                        
                        st.success(f"✅ Processed {len(selected_files) + (1 if youtube_url else 0)} items")
                        st.rerun()
                        
                    except Exception as e:
                        st.error(f"Error processing documents: {e}")
            else:
                st.warning("Please select course files or provide a YouTube URL")
        
        # Show processed files
        if st.session_state.uploaded_files:
            st.subheader("📚 Processed Content")
            for i, file_info in enumerate(st.session_state.uploaded_files):
                file_type_display = {
                    'course_document': '📄 Course Document',
                    'video': '🎥 Video',
                    'document': '📄 Document'
                }.get(file_info['type'], file_info['type'].title())
                
                # Create a styled expander with better visual appeal
                expander_style = """
                <style>
                .streamlit-expanderHeader {
                    background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                    color: white;
                    border-radius: 10px;
                    padding: 10px;
                    margin: 5px 0;
                    box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
                }
                .streamlit-expanderContent {
                    background: #f8f9fa;
                    border-radius: 0 0 10px 10px;
                    padding: 15px;
                }
                </style>
                """
                st.markdown(expander_style, unsafe_allow_html=True)
                
                with st.expander(f"📄 {file_type_display}: {file_info['name']}", expanded=False):
                    # Enhanced content preview with prettified formatting
                    st.markdown("### 📄 Content Preview")
                    
                    # Create a beautiful container for the content preview
                    st.markdown("""
                    <div style="
                        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                        padding: 20px;
                        border-radius: 15px;
                        margin: 10px 0;
                        box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1);
                        border: 1px solid rgba(255, 255, 255, 0.2);
                    ">
                        <div style="
                            background: white;
                            padding: 20px;
                            border-radius: 10px;
                            box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
                            max-height: 400px;
                            overflow-y: auto;
                        ">
                    """, unsafe_allow_html=True)
                    
                    # Format the content preview with enhanced styling
                    preview_text = file_info['text'][:2000] + "..." if len(file_info['text']) > 2000 else file_info['text']
                    formatted_preview = self._format_content_preview(preview_text, file_info['name'])
                    st.markdown(formatted_preview)
                    
                    st.markdown("""
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    if st.button(f"🗑️ Remove", key=f"remove_{i}"):
                        st.session_state.uploaded_files.pop(i)
                        st.rerun()
    
    def render_ai_notes_generation(self):
        """Render AI notes generation interface - exact same as LMS-AI-Assistant"""
        st.subheader("Generate Lecture-Style Notes")
        
        # Check for local files in downloads directory
        downloads_dir = Path("downloads")
        local_files = []
        
        if downloads_dir.exists():
            # Find all supported files in downloads directory
            supported_extensions = ['.pdf', '.docx', '.doc', '.txt', '.pptx', '.ppsx']
            for ext in supported_extensions:
                local_files.extend(downloads_dir.rglob(f'*{ext}'))
        
        if not local_files and not st.session_state.uploaded_files:
            st.info("📁 **No Local Files Found**")
            st.markdown("""
            **To generate AI notes, you need to:**
            1. Download files locally using the Download section
            2. Or upload files manually using the file uploader below
            
            This approach works around Canvas API permission restrictions by processing files locally.
            """)
            
            # File uploader as fallback
            uploaded_files = st.file_uploader(
                "Upload files for AI processing",
                type=['pdf', 'docx', 'doc', 'txt', 'pptx', 'ppsx'],
                accept_multiple_files=True,
                help="Upload files if you don't have any in the downloads folder"
            )
            
            if uploaded_files:
                st.session_state.uploaded_files = []
                for uploaded_file in uploaded_files:
                    # Save uploaded file temporarily
                    temp_path = Path("temp") / uploaded_file.name
                    temp_path.parent.mkdir(exist_ok=True)
                    with open(temp_path, "wb") as f:
                        f.write(uploaded_file.getbuffer())
                    
                    # Extract text
                    result = self._extract_text_from_file(temp_path)
                    if result.get('success'):
                        st.session_state.uploaded_files.append({
                            'name': uploaded_file.name,
                            'text': result.get('full_text', ''),
                            'type': 'uploaded_document',
                            'file_path': str(temp_path)
                        })
                    temp_path.unlink()  # Clean up temp file
                
                if st.session_state.uploaded_files:
                    st.success(f"✅ Processed {len(st.session_state.uploaded_files)} uploaded files!")
        
        # Show available files
        all_files = []
        if local_files:
            for file_path in local_files[:10]:  # Limit to first 10 files
                result = self._extract_text_from_file(file_path)
                if result.get('success'):
                    all_files.append({
                        'name': file_path.name,
                        'text': result.get('full_text', ''),
                        'type': 'local_document',
                        'file_path': str(file_path)
                    })
        
        if st.session_state.uploaded_files:
            all_files.extend(st.session_state.uploaded_files)
        
        if not all_files:
            return
        
        # Exact same UI as LMS-AI-Assistant
        coln1, coln2 = st.columns([3, 1])
        with coln1:
            custom_title = st.text_input("Notes Title (optional)")
        with coln2:
            notes_chunk_size = st.number_input("Chunk size", min_value=400, max_value=4000, value=1200, step=100)
        
        if st.button("Start Note-Making"):
            with st.spinner("Generating notes..."):
                # Gather all text from files
                texts = [f['text'] for f in all_files if f.get('text')]
                
                if not texts:
                    st.warning("No content found in files to generate notes from.")
                    return
                
                st.markdown("### Notes (Generating)")
                placeholder = st.empty()
                col_a, col_b = st.columns([3,1])
                with col_b:
                    group_size = st.number_input("Chunks/group", min_value=1, max_value=10, value=3)

                sections = []
                for idx, sec in enumerate(self.iter_generate_notes_from_texts(texts, title=custom_title, group_size=int(group_size)), 1):
                    sections.append(sec)
                    with placeholder.container():
                        for i, s in enumerate(sections, 1):
                            with st.expander(f"Section {i}", expanded=(i == idx)):
                                st.markdown(s)
                st.success("Notes generated.")
                
                # Store generated notes in session state
                st.session_state.current_notes = "\n\n".join(sections)
                st.session_state.notes_generated = True
    
    def render_qa_chatbot(self):
        """Render QA chatbot interface"""
        st.subheader("💬 QA Chatbot")
        
        if not st.session_state.uploaded_files:
            st.info("Upload documents first to use the chatbot")
            return
        
        # Chat interface
        for message in st.session_state.chat_history:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])
        
        # Chat input
        if prompt := st.chat_input("Ask a question about your documents..."):
            # Add user message
            st.session_state.chat_history.append({"role": "user", "content": prompt})
            
            with st.chat_message("user"):
                st.markdown(prompt)
            
            # Generate response
            with st.chat_message("assistant"):
                with st.spinner("Thinking..."):
                    try:
                        # Combine uploaded content for context
                        context = "\n\n".join([f['text'] for f in st.session_state.uploaded_files])
                        
                        # Generate response using semantic search and QA
                        response = answer_question_with_context(prompt, context)
                        st.markdown(response)
                        
                        # Add assistant response
                        st.session_state.chat_history.append({"role": "assistant", "content": response})
                        
                    except Exception as e:
                        error_msg = f"Sorry, I encountered an error: {e}"
                        st.markdown(error_msg)
                        st.session_state.chat_history.append({"role": "assistant", "content": error_msg})
        
        # Clear chat button
        if st.button("🗑️ Clear Chat"):
            st.session_state.chat_history = []
            st.rerun()
    
    def render_rag_resources_tab(self):
        """Render recommended resources based on extracted topics"""
        st.subheader("🔗 Recommended Resources")
        
        if not st.session_state.extracted_topics:
            st.info("Generate AI notes first to see recommended resources")
            return
        
        st.markdown("Based on your uploaded content, here are some recommended resources:")
        
        for topic in st.session_state.extracted_topics[:5]:  # Show top 5 topics
            with st.expander(f"📚 {topic}"):
                col1, col2 = st.columns(2)
                
                with col1:
                    st.markdown("**📺 Recommended Videos**")
                    # This would integrate with YouTube API or use predefined resources
                    st.markdown(f"• [Introduction to {topic}](https://youtube.com/search?q={topic.replace(' ', '+')})")
                    st.markdown(f"• [Advanced {topic}](https://youtube.com/search?q=advanced+{topic.replace(' ', '+')})")
                
                with col2:
                    st.markdown("**📄 Recommended Articles**")
                    # This would integrate with academic databases or use predefined resources
                    st.markdown(f"• [Wikipedia: {topic}](https://en.wikipedia.org/wiki/{topic.replace(' ', '_')})")
                    st.markdown(f"• [Khan Academy: {topic}](https://www.khanacademy.org/search?search_again=1&page_search_query={topic.replace(' ', '+')})")
        
        if st.button("🔍 Explain connections between topics"):
            with st.spinner("Analyzing topic relationships..."):
                try:
                    # Generate explanation of how topics connect
                    topics_text = ", ".join(st.session_state.extracted_topics[:5])
                    explanation = f"The main topics in your content ({topics_text}) are interconnected through various academic and practical relationships. Understanding these connections can help you build a more comprehensive knowledge framework."
                    st.info(explanation)
                except Exception as e:
                    st.error(f"Error analyzing connections: {e}")
    
    def run(self):
        """Main app runner"""
        self.initialize_session_state()
        self.render_header()
        
        # Sidebar authentication
        is_authenticated = self.render_login_section()
        
        if is_authenticated:
            # Main content
            self.render_course_selection()
            
            if st.session_state.selected_course:
                self.render_course_overview()
                # self.render_workflow_status()  # Commented out to hide workflow dashboard
                
                if RAG_AVAILABLE:
                    tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs(["📝 Assignments", "📚 Modules", "📁 Files", "🤖 AI Notes", "🎓 Guided Learning", "💬 RAG Demo"])
                else:
                    tab1, tab2, tab3, tab4, tab5 = st.tabs(["📝 Assignments", "📚 Modules", "📁 Files", "🤖 AI Notes", "🎓 Guided Learning"])
                
                with tab1:
                    self.render_assignments_section()
                
                with tab2:
                    self.render_modules_section()
                
                with tab3:
                    self.render_files_section()
                
                with tab4:
                    if RAG_AVAILABLE:
                        self.render_ai_notes_section()
                    else:
                        st.error("🤖 RAG Demo components not available. Please check dependencies.")
                        st.info("Install required packages: pip install sentence-transformers chromadb")
                
                with tab5:
                    self.render_guided_learning_section()
                
                if RAG_AVAILABLE:
                    with tab6:
                        self.render_rag_demo_section()
        else:
            # Show instructions when not authenticated
            st.markdown("""
            ## Welcome to Canvas LMS Course Explorer! 🎉
            
            This application allows you to:
            - 🔐 Securely connect to your Canvas LMS account
            - 📚 Browse and explore your courses
            - 📝 View assignments and due dates
            - 📚 Explore course modules and content
            - 📊 Export course data for analysis
            
            ### Getting Started:
            1. **Login**: Use the sidebar to enter your Canvas URL and API token
            2. **Select Course**: Choose from your available courses
            3. **Explore Data**: View assignments, modules, and other course information
            4. **Export**: Download course data in various formats
            
            ### Note:
            This app focuses on data you have access to as a student. Some features like file downloads may not be available due to Canvas permissions.
            """)
    
    def render_guided_learning_section(self):
        """Render the guided learning section with step-by-step learning"""
        st.header("🎓 Guided Learning Assistant")
        st.markdown("Learn any topic step-by-step with personalized content, YouTube videos, and articles!")
        
        # Initialize session state for guided learning
        if 'learning_messages' not in st.session_state:
            st.session_state.learning_messages = []
        if 'learning_initialized' not in st.session_state:
            st.session_state.learning_initialized = False
        if 'learning_cooldown' not in st.session_state:
            st.session_state.learning_cooldown = 0
        
        # Get notes sections for context
        notes_sections = []
        try:
            # Try to get notes from AI notes cache first
            if hasattr(st.session_state, 'ai_notes_cache') and st.session_state.ai_notes_cache:
                for file_path, notes_data in st.session_state.ai_notes_cache.items():
                    if isinstance(notes_data, dict) and 'notes' in notes_data:
                        notes_sections.append(notes_data['notes'])
            
            # If no notes available, show message
            if not notes_sections:
                st.info("💡 Generate AI notes first to enable guided learning with course content context!")
                st.markdown("### Quick Start Learning")
                st.markdown("You can still use guided learning without course notes. Just ask me about any topic!")
        except Exception as e:
            st.warning(f"Could not load course notes: {e}")
        
        # Learning interface
        col1, col2 = st.columns([2, 1])
        
        with col1:
            st.subheader("💬 Learning Chat")
            
            # Display conversation history
            for message in st.session_state.learning_messages:
                with st.chat_message(message["role"]):
                    st.markdown(message["content"])
            
            # Chat input
            user_input = st.chat_input("Ask me about any topic you want to learn...")
            
            if user_input:
                import time
                now = time.time()
                
                # Cooldown check
                if now - st.session_state.learning_cooldown < 3.0:
                    st.warning("Please wait a moment before sending another message.")
                    st.stop()
                
                with st.spinner("Learning assistant is thinking..."):
                    try:
                        from utils.learning_mode import adaptive_learning_response
                        
                        # Get adaptive response
                        adaptive_response = adaptive_learning_response(
                            user_input, 
                            notes_sections, 
                            st.session_state.learning_messages, 
                            profile_id="canvas_learner"
                        )
                        
                        reply = adaptive_response["response"]
                        resources = adaptive_response["resources"]
                        quiz_content = adaptive_response.get("quiz", "")
                        learning_prefs = adaptive_response.get("learning_preferences", {})
                        
                        # Add user message
                        st.session_state.learning_messages.append({"role": "user", "content": user_input})
                        
                        # Add assistant response
                        st.session_state.learning_messages.append({"role": "assistant", "content": reply})
                        
                        # Update cooldown
                        st.session_state.learning_cooldown = now
                        
                        # Display learning preferences if detected
                        if learning_prefs and any(learning_prefs.values()):
                            with st.expander("🎯 Detected Learning Preferences", expanded=False):
                                for key, value in learning_prefs.items():
                                    if value:
                                        st.write(f"**{key.replace('_', ' ').title()}**: {value}")
                        
                        # Display resources
                        if resources:
                            self.display_learning_resources(resources)
                        
                        # Display quiz if available
                        if quiz_content:
                            with st.expander("📝 Quick Quiz", expanded=False):
                                st.markdown(quiz_content)
                        
                        st.rerun()
                        
                    except Exception as e:
                        st.error(f"Error in learning assistant: {e}")
                        st.session_state.learning_messages.append({"role": "user", "content": user_input})
                        st.session_state.learning_messages.append({"role": "assistant", "content": "I'm sorry, I encountered an error. Please try again."})
                        st.rerun()
            
        with col2:
            st.subheader("📚 Learning Resources")
            
            # Topic-based recommendations
            if notes_sections:
                try:
                    from utils.learning_mode import extract_topics_from_notes
                    from utils.web_search import recommend_articles_ddg, recommend_youtube_ddg
                    
                    topics = extract_topics_from_notes(notes_sections)
                    
                    if topics:
                        st.markdown("### Course Topics")
                        for topic in topics[:5]:  # Show top 5 topics
                            with st.expander(f"📖 {topic}", expanded=False):
                                # Get videos and articles for this topic
                                vlinks = recommend_youtube_ddg(topic, limit=2)
                                alinks = recommend_articles_ddg(topic, limit=2)
                                
                                if vlinks:
                                    st.markdown("**🎥 Videos:**")
                                    for link in vlinks:
                                        st.markdown(f"- [{link['title']}]({link['url']})")
                                else:
                                    st.write("No videos found for this topic.")
                                
                                if alinks:
                                    st.markdown("**📄 Articles:**")
                                    for link in alinks:
                                        st.markdown(f"- [{link['title']}]({link['url']})")
                                else:
                                    st.write("No articles found for this topic.")
                    else:
                        st.info("No topics found in course notes.")
                        
                except Exception as e:
                    st.warning(f"Could not load topic recommendations: {e}")
            else:
                st.info("Generate AI notes to see topic-based recommendations!")
            
            # Learning suggestions
            st.markdown("### 💡 Learning Tips")
            st.markdown("""
            - **Be specific**: Ask about particular concepts or topics
            - **Mention preferences**: Say "I prefer videos" or "I like hands-on practice"
            - **Ask for examples**: Request practical examples or use cases
            - **Request quizzes**: Ask for practice questions to test understanding
            """)
    
    def display_learning_resources(self, resources: Dict[str, Any]):
        """Display learning resources in an organized way"""
        if not any(resources.values()):
            return
        
        st.markdown("### 📚 Recommended Resources")
        
        # Videos
        if resources.get("videos"):
            with st.expander("🎥 Video Tutorials", expanded=True):
                for video in resources["videos"]:
                    st.markdown(f"- [{video['title']}]({video['url']})")
        
        # Articles
        if resources.get("articles"):
            with st.expander("📄 Articles & Documentation", expanded=True):
                for article in resources["articles"]:
                    st.markdown(f"- [{article['title']}]({article['url']})")
        
        # Assignments
        if resources.get("assignments"):
            with st.expander("📝 Practice Assignments", expanded=True):
                for assignment in resources["assignments"]:
                    st.markdown(f"**{assignment['title']}**")
                    st.markdown(assignment['content'])
        
        # Projects
        if resources.get("projects"):
            with st.expander("🚀 Project Ideas", expanded=True):
                for project in resources["projects"]:
                    st.markdown(f"**{project['title']}**")
                    st.markdown(project['content'])


def main():
    """Main function"""
    app = CanvasCourseExplorer()
    app.run()

if __name__ == "__main__":
    main()
