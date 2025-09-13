"""PowerPoint Text Extractor
Extracts text from PowerPoint presentations (.pptx, .ppsx)
"""

import os
import logging
from typing import Dict, Any, List
from datetime import datetime

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. PowerPoint extraction will not work.")

logger = logging.getLogger(__name__)

class PPTExtractor:
    """Extract text from PowerPoint presentations"""
    
    def __init__(self):
        """Initialize the PowerPoint extractor"""
        self.supported_extensions = ['.pptx', '.ppsx']
        
        if not PPTX_AVAILABLE:
            logger.warning("PowerPoint extraction requires python-pptx package")
    
    def extract_text(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text from a PowerPoint file
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            Dictionary containing extraction results
        """
        if not PPTX_AVAILABLE:
            return {
                'success': False,
                'error': 'python-pptx package not available. Install with: pip install python-pptx',
                'file_path': file_path
            }
        
        if not os.path.exists(file_path):
            return {
                'success': False,
                'error': f'File not found: {file_path}',
                'file_path': file_path
            }
        
        try:
            # Load the presentation
            prs = Presentation(file_path)
            
            # Extract text from all slides
            slides_text = []
            slide_count = 0
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    slides_text.append({
                        'slide_number': slide_num,
                        'text': '\n'.join(slide_text)
                    })
                    slide_count += 1
            
            # Combine all slide text
            full_text = '\n\n'.join([slide['text'] for slide in slides_text])
            
            # Get file metadata
            file_stats = os.stat(file_path)
            file_size = file_stats.st_size
            
            return {
                'success': True,
                'text': full_text,
                'slides': slides_text,
                'metadata': {
                    'file_path': file_path,
                    'file_name': os.path.basename(file_path),
                    'file_size': file_size,
                    'slide_count': slide_count,
                    'total_slides': len(prs.slides),
                    'extraction_method': 'python-pptx',
                    'extracted_at': datetime.now().isoformat()
                }
            }
            
        except Exception as e:
            logger.error(f"Error extracting text from PowerPoint {file_path}: {e}")
            return {
                'success': False,
                'error': str(e),
                'file_path': file_path
            }
    
    def extract_slide_notes(self, file_path: str) -> Dict[str, Any]:
        """
        Extract speaker notes from PowerPoint slides
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            Dictionary containing slide notes
        """
        if not PPTX_AVAILABLE:
            return {
                'success': False,
                'error': 'python-pptx package not available',
                'file_path': file_path
            }
        
        try:
            prs = Presentation(file_path)
            notes_data = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        notes_data.append({
                            'slide_number': slide_num,
                            'notes': notes_text
                        })
            
            return {
                'success': True,
                'notes': notes_data,
                'total_slides_with_notes': len(notes_data)
            }
            
        except Exception as e:
            logger.error(f"Error extracting notes from PowerPoint {file_path}: {e}")
            return {
                'success': False,
                'error': str(e),
                'file_path': file_path
            }
    
    def get_presentation_info(self, file_path: str) -> Dict[str, Any]:
        """
        Get basic information about the PowerPoint presentation
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            Dictionary containing presentation information
        """
        if not PPTX_AVAILABLE:
            return {
                'success': False,
                'error': 'python-pptx package not available'
            }
        
        try:
            prs = Presentation(file_path)
            
            # Get core properties if available
            core_props = prs.core_properties
            
            return {
                'success': True,
                'info': {
                    'slide_count': len(prs.slides),
                    'title': getattr(core_props, 'title', None),
                    'author': getattr(core_props, 'author', None),
                    'subject': getattr(core_props, 'subject', None),
                    'created': getattr(core_props, 'created', None),
                    'modified': getattr(core_props, 'modified', None)
                }
            }
            
        except Exception as e:
            logger.error(f"Error getting PowerPoint info {file_path}: {e}")
            return {
                'success': False,
                'error': str(e)
            }